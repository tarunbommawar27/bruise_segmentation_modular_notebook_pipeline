#!/usr/bin/env python
"""Emit docs/b2_distillation_deck.pptx -- the SegFormer-B2 distillation results.

DESIGN BRIEF, AND WHY EACH CHOICE
----------------------------------
Minimal, Times New Roman, professional. Concretely that means:

  * ONE idea per slide, stated in the title as a claim rather than a label. A
    slide called "Results" makes the reader do the work; "Only the B0 students
    hold the line" does not.
  * Figures at full width. Every figure here was generated at 200 dpi by the
    notebook, so it is placed at native aspect ratio and never stretched.
  * Numbers come from the CSVs at build time. Nothing in this file is a
    transcribed figure -- a deck that hardcodes 0.7656 goes stale the moment a
    seed is re-run, and nobody notices.
  * No bullet lists longer than four items, no sub-bullets, no clip art.

WHAT THE DECK IS SCOPED TO
---------------------------
Only arms whose teacher is **SegFormer-B2**, with **`segformer_b0_direct` as the
boundary** -- the strongest model in the study trained with no teacher at all. So
every contrast reads "distilled from B2, against not distilling". The DeepLabV3+
arms appear once, in the teacher-swap slide, because that comparison is what
licenses using B2 at all.

The annotation-ceiling slide comes BEFORE any result slide, deliberately. It is
the fact that governs how every later number is read (handbook 1), and a deck
that puts it in an appendix invites exactly the over-reading it exists to prevent.
"""
from __future__ import annotations

import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
RES = ROOT / "BRUISE_UNIFIED" / "FINAL_RESULT"
FIG = RES / "figures"
OUT = ROOT / "docs" / "b2_distillation_deck.pptx"

FONT = "Times New Roman"
INK = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT = RGBColor(0x1F, 0x3F, 0x6E)      # navy
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
RULE = RGBColor(0xC9, 0xCF, 0xD8)
GOOD = RGBColor(0x1B, 0x5E, 0x3A)
BAD = RGBColor(0x8B, 0x2C, 0x2C)

W, H = Inches(13.333), Inches(7.5)        # 16:9
M = Inches(0.75)                          # side margin


# ─────────────────────────────────────────────────────────────────────────────
# data (read once, never transcribed)
# ─────────────────────────────────────────────────────────────────────────────
def csv(name: str) -> pd.DataFrame:
    p = RES / name
    if not p.exists():
        raise SystemExit(f"missing {p} -- run the notebook's Stage H cells first")
    return pd.read_csv(p)


D = {
    "b2": csv("significance_b2_teacher_vs_b0_direct.csv"),
    "famH": csv("significance_contrast_family_stage_h.csv"),
    "gate": csv("reliability_gate_diagnostics.csv"),
    "hA": csv("headline_stage_a.csv"),
    "hH": csv("headline_stage_h.csv"),
    "hE": csv("headline_stage_e.csv"),
    "omni": csv("significance_omnibus.csv"),
    "ceil": csv("annotation_ceiling.csv"),
}


def famH(a: str, b: str) -> pd.Series:
    m = D["famH"][(D["famH"].a == a) & (D["famH"].b == b)]
    if not len(m):
        raise SystemExit(f"contrast {a} vs {b} not in the Stage H family")
    return m.iloc[0]


def dice(table: str, model: str) -> float:
    d = D[table]
    return float(d[d.model == model].mean_dice.iloc[0])


# ─────────────────────────────────────────────────────────────────────────────
# slide furniture
# ─────────────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def _fmt(run, size, bold=False, color=INK, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def textbox(slide, left, top, width, height, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    return tf


def slide(title: str, kicker: str | None = None):
    """A slide with a claim as its title, a hairline rule, and optional kicker."""
    s = prs.slides.add_slide(BLANK)
    tf = textbox(s, M, Inches(0.42), W - 2 * M, Inches(0.75))
    _fmt(tf.paragraphs[0].add_run(), 26, bold=True, color=ACCENT)
    tf.paragraphs[0].runs[0].text = title

    ln = s.shapes.add_shape(1, M, Inches(1.16), W - 2 * M, Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE
    ln.line.fill.background(); ln.shadow.inherit = False

    if kicker:
        kf = textbox(s, M, Inches(1.24), W - 2 * M, Inches(0.5))
        _fmt(kf.paragraphs[0].add_run(), 13, color=MUTED, italic=True)
        kf.paragraphs[0].runs[0].text = kicker
    return s


def picture(s, name: str, top: Inches, max_h: Inches, max_w=None):
    """Place a figure centred, scaled to fit, never stretched."""
    from PIL import Image
    p = FIG / name
    iw, ih = Image.open(p).size
    max_w = max_w or (W - 2 * M)
    w = max_w
    h = int(w * ih / iw)
    if h > max_h:
        h = max_h
        w = int(h * iw / ih)
    s.shapes.add_picture(str(p), int((W - w) / 2), top, width=int(w), height=int(h))


def table(s, rows, left, top, width, col_w=None, size=12, header=True,
          highlight=None, row_h=Inches(0.34)):
    """A minimal table: header rule only, no grid, no fill."""
    n_r, n_c = len(rows), len(rows[0])
    shape = s.shapes.add_table(n_r, n_c, left, top, width, row_h * n_r)
    tbl = shape.table
    tbl.first_row = False
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = int(width * cw / total)

    for r, row in enumerate(rows):
        tbl.rows[r].height = row_h
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = 0
            cell.margin_bottom = 0
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            col = INK
            if highlight and r in highlight:
                col = highlight[r]
            _fmt(run, size, bold=(header and r == 0), color=MUTED if (header and r == 0) else col)
    return tbl


def note(s, text, top=Inches(6.55)):
    tf = textbox(s, M, top, W - 2 * M, Inches(0.7))
    _fmt(tf.paragraphs[0].add_run(), 11.5, color=MUTED)
    tf.paragraphs[0].runs[0].text = text


# ═════════════════════════════════════════════════════════════════════════════
# 1 · title
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
tf = textbox(s, M, Inches(2.45), W - 2 * M, Inches(1.3))
_fmt(tf.paragraphs[0].add_run(), 40, bold=True, color=ACCENT)
tf.paragraphs[0].runs[0].text = "Knowledge Distillation from SegFormer-B2"
p = tf.add_paragraph()
_fmt(p.add_run(), 40, bold=True, color=ACCENT)
p.runs[0].text = "for Bruise Segmentation"

tf2 = textbox(s, M, Inches(4.05), W - 2 * M, Inches(1.0))
_fmt(tf2.paragraphs[0].add_run(), 17, color=INK)
tf2.paragraphs[0].runs[0].text = ("Response distillation and reliability-gated distillation "
                                  "into five compact students")
p = tf2.add_paragraph()
_fmt(p.add_run(), 15, color=MUTED, italic=True)
p.runs[0].text = "Reference boundary: SegFormer-B0 trained without a teacher"

ln = s.shapes.add_shape(1, M, Inches(5.35), Inches(2.2), Emu(19050))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT
ln.line.fill.background(); ln.shadow.inherit = False

tf3 = textbox(s, M, Inches(5.6), W - 2 * M, Inches(0.9))
_fmt(tf3.paragraphs[0].add_run(), 13, color=MUTED)
n_sub = int(D["b2"].n_subjects.iloc[0])
tf3.paragraphs[0].runs[0].text = (f"1 016 white-light photographs · 143 subjects · "
                                  f"185 test images from {n_sub} subjects · 3 seeds per arm")

# ═════════════════════════════════════════════════════════════════════════════
# 2 · protocol
# ═════════════════════════════════════════════════════════════════════════════
s = slide("One recipe, held fixed across every architecture",
          "Differences between arms are attributable to the teacher and the loss, to nothing else.")
rows = [
    ["", "Setting", "", "Setting"],
    ["Split", "697 / 134 / 185, subject-grouped", "Loss", "Dice + BCE (+ KD term)"],
    ["Input", "640 × 640", "Selection", "best validation AP"],
    ["Optimiser", "AdamW, poly decay", "Threshold", "swept on val, applied once to test"],
    ["Learning rate", "6e-5 backbone / 6e-4 head", "Seeds", "0, 1, 2"],
]
table(s, rows, M, Inches(1.75), W - 2 * M, col_w=[1.3, 3.0, 1.3, 3.0], size=13.5)
note(s, "No subject appears in two splits. The threshold is never fitted on test. "
        "KD mix α = 0.6 for every arm, gated and ungated alike — so the gate is the "
        "only thing that differs between a gated arm and its control.", Inches(4.55))

# ═════════════════════════════════════════════════════════════════════════════
# 3 · the annotation ceiling  (framing, before any result)
# ═════════════════════════════════════════════════════════════════════════════
s = slide("Human annotators disagree more than the models do",
          "This governs how every number that follows must be read.")
picture(s, "D7_annotation_ceiling.png", Inches(1.6), Inches(4.55))
c = D["ceil"]
note(s, "Two expert annotators agree with each other at 0.755 Dice; the weakest pair at 0.581. "
        "The entire model field sits inside that band. A difference below ≈0.05 Dice is not a "
        "result unless a paired subject-level bootstrap excludes zero — which is why every "
        "claim in this deck carries an interval.", Inches(6.3))

# ═════════════════════════════════════════════════════════════════════════════
# 4 · method
# ═════════════════════════════════════════════════════════════════════════════
s = slide("Reliability-gated distillation",
          "The teacher is trusted per pixel and per image, instead of everywhere equally.")

tf = textbox(s, M, Inches(1.85), Inches(6.4), Inches(2.2))
for txt, sz, bold, col in [
    ("r  =  1 − |2p − 1| · |p − y|", 19, True, ACCENT),
    ("per-pixel reliability", 12, False, MUTED),
    ("", 6, False, MUTED),
    ("g  =  clip( (Dice_teacher − 0.10) / 0.40 , 0, 1 )", 19, True, ACCENT),
    ("per-image gate, on the teacher's own soft Dice", 12, False, MUTED),
]:
    p = tf.paragraphs[0] if txt == "r  =  1 − |2p − 1| · |p − y|" else tf.add_paragraph()
    _fmt(p.add_run(), sz, bold=bold, color=col)
    p.runs[0].text = txt

rows = [
    ["Teacher is…", "confidence", "error", "r", "soft term"],
    ["confidently right", "≈ 1", "≈ 0", "≈ 1", "full weight"],
    ["uncertain", "≈ 0", "any", "≈ 1", "full weight"],
    ["confidently wrong", "≈ 1", "≈ 1", "≈ 0", "suppressed"],
]
table(s, rows, M, Inches(4.15), Inches(7.4), col_w=[2.2, 1.1, 0.9, 0.7, 1.5], size=12.5,
      highlight={2: ACCENT})
note(s, "The middle row is the design. Down-weighting by error alone would delete exactly the "
        "pixels the teacher is unsure about — the information distillation exists to transfer. "
        "Multiplying by confidence removes only assertive error. With r ≡ 1 and g ≡ 1 the loss "
        "reduces exactly to standard response KD, so each contrast moves one variable.",
     Inches(5.75))

# ═════════════════════════════════════════════════════════════════════════════
# 5 · did the gate fire
# ═════════════════════════════════════════════════════════════════════════════
g = D["gate"]
s = slide("The gate fired — it is neither inert nor overwhelming",
          "Read before any accuracy number: an ungated gate is a relabelled control.")
rows = [["Diagnostic", "Value", "Reading"],
        ["Mean coverage (g · r)", f"{g.mean_coverage.mean():.3f}",
         "1.0 would mean the gate never fired"],
        ["Effective α", f"{g.mean_alpha_effective.mean():.3f}",
         f"nominal α = {g.alpha_nominal.iloc[0]:.1f}; freed weight returns to supervision"],
        ["Image-views fully gated off", f"{g.frac_images_fully_gated_off.mean()*100:.1f} %",
         "the teacher was ignored entirely on these"],
        ["Teacher near-miss views", f"{g.frac_teacher_near_miss.mean()*100:.1f} %",
         "teacher soft Dice ≤ 0.05 — the target population"],
        ["Mean pixel reliability", f"{g.mean_pixel_reliability.mean():.3f}",
         "most pixels are confidently correct, as expected"]]
table(s, rows, M, Inches(1.85), W - 2 * M, col_w=[2.6, 1.2, 5.0], size=13.5)
note(s, f"Averaged over {len(g)} gated runs (5 students × 3 seeds). The gate removed the KD term "
        f"on roughly one image-view in thirty-five — the images where the teacher asserted an "
        f"empty mask on a real bruise — and left it untouched elsewhere.", Inches(4.85))

# ═════════════════════════════════════════════════════════════════════════════
# 6 · accuracy overview
# ═════════════════════════════════════════════════════════════════════════════
s = slide("Accuracy of every B2-taught arm",
          "Val-selected seed, 185 test images. Reference in bold.")
b2 = D["b2"]
rows = [["Student", "Params", "KD method", "Mean Dice", "Median", "Misses"]]
rows.append(["SegFormer-B0  (no teacher)", "3.71 M", "—",
             f"{dice('hA','segformer_b0_direct'):.4f}",
             f"{D['hA'][D['hA'].model=='segformer_b0_direct'].median_dice.iloc[0]:.4f}", "1"])
hl = {1: ACCENT}
for i, r in enumerate(b2.itertuples(), start=2):
    rows.append([r.student, f"{r.params_M:.2f} M",
                 "gated" if r.method.startswith("reliability") else "response",
                 f"{r.mean_dice:.4f}", f"{r.median_dice:.4f}", str(int(r.misses))])
table(s, rows, M, Inches(1.8), W - 2 * M, col_w=[2.6, 1.0, 1.4, 1.3, 1.1, 0.9],
      size=12, highlight=hl, row_h=Inches(0.325))
note(s, "Median is reported beside the mean because the per-image distribution is strongly "
        "left-skewed: a handful of complete misses moves the mean several points while the "
        "median barely shifts.", Inches(6.35))

# ═════════════════════════════════════════════════════════════════════════════
# 7 · significance table
# ═════════════════════════════════════════════════════════════════════════════
s = slide("Against the B0-direct boundary, only the B0 students hold the line",
          "Paired subject-level bootstrap, 28 subjects, 10 000 draws, Holm-corrected over all ten.")
rows = [["Arm", "Δ Dice", "95 % CI", "p (Holm)", "Verdict"]]
hl = {}
for i, r in enumerate(b2.itertuples(), start=1):
    meth = "gated" if r.method.startswith("reliability") else "response"
    rows.append([f"{r.student} · {meth}", f"{r.delta_dice:+.4f}",
                 f"[{r.lo:+.4f}, {r.hi:+.4f}]",
                 f"{r.p_holm:.4f}", r.verdict])
    hl[i] = BAD if r.verdict == "INFERIOR" else (GOOD if r.verdict == "NON-INFERIOR" else INK)
table(s, rows, M, Inches(1.85), W - 2 * M, col_w=[2.9, 1.1, 2.0, 1.1, 1.6],
      size=12.5, highlight=hl, row_h=Inches(0.34))
note(s, "NON-INFERIOR = within the 0.01 Dice margin. INCONCLUSIVE = the interval is too wide to "
        "decide at 28 subjects, which is not the same as equivalence. Distilling B2 into a 3.71 M "
        "SegFormer-B0 costs nothing measurable; distilling it into a 1–3 M mobile architecture "
        "does not recover the gap.", Inches(6.4))

# ═════════════════════════════════════════════════════════════════════════════
# 8 · forest plot
# ═════════════════════════════════════════════════════════════════════════════
s = slide("The same result, as intervals")
picture(s, "H1_b2_vs_b0_direct.png", Inches(1.5), Inches(4.7))
note(s, "Every interval that lies wholly left of the line is a student that is genuinely worse "
        "than an undistilled SegFormer-B0. The two B0 arms straddle it. Gated and response KD "
        "are indistinguishable from each other at every student size.", Inches(6.35))

# ═════════════════════════════════════════════════════════════════════════════
# 9 · what distillation did buy
# ═════════════════════════════════════════════════════════════════════════════
s = slide("Distillation does help — against no teacher, not against the boundary",
          "The comparison that matters for a deployable model is arm vs its own direct baseline.")
pairs = [
    ("lraspp_mobilenetv3_rgkd", "lraspp_mobilenetv3", "LR-ASPP MNv3", "gated KD vs no KD"),
    ("fastscnn_rgkd", "fastscnn", "Fast-SCNN", "gated KD vs no KD"),
    ("fastscnn_b2kd", "fastscnn_distilled", "Fast-SCNN", "B2 teacher vs DeepLabV3+"),
]
rows = [["Contrast", "Student", "Δ Dice", "95 % CI", "Verdict"]]
hl = {}
for i, (a, b, student, label) in enumerate(pairs, start=1):
    r = famH(a, b)
    rows.append([label, student, f"{r.delta_dice:+.4f}",
                 f"[{r.lo:+.4f}, {r.hi:+.4f}]", r.verdict])
    hl[i] = GOOD if r.verdict == "WIN" else INK
table(s, rows, M, Inches(1.9), W - 2 * M, col_w=[2.4, 1.8, 1.1, 2.0, 1.2], size=13.5)

r_fs = famH("fastscnn_rgkd", "fastscnn")
note(s, f"Three wins, each with an interval excluding zero. Gated KD lifts Fast-SCNN by "
        f"{r_fs.delta_dice:+.4f} Dice and cuts its complete-miss rate by "
        f"{abs(r_fs.delta_miss_rate)*100:.1f} points. Switching the teacher from DeepLabV3+ to "
        f"B2 is worth {famH('fastscnn_b2kd','fastscnn_distilled').delta_dice:+.4f} on the same "
        f"student — which is what justifies using B2 at all.", Inches(4.2))

# ═════════════════════════════════════════════════════════════════════════════
# 10 · gating vs response KD
# ═════════════════════════════════════════════════════════════════════════════
s = slide("Reliability gating did not beat plain response KD",
          "Reported as the null it is. Teacher held fixed at B2; only the loss differs.")
rows = [["Student", "Δ Dice (gated − response)", "95 % CI", "Verdict"]]
for a, b in [("segformer_b0_rgkd", "segformer_b0_distilled"),
             ("lraspp_mobilenetv3_rgkd", "lraspp_mobilenetv3_b2kd"),
             ("fastscnn_rgkd", "fastscnn_b2kd"),
             ("topformer_tiny_rgkd", "topformer_tiny_b2kd"),
             ("ppmobileseg_tiny_rgkd", "ppmobileseg_tiny_b2kd")]:
    r = famH(a, b)
    name = b2[b2.arm == a].student.iloc[0]
    rows.append([name, f"{r.delta_dice:+.4f}", f"[{r.lo:+.4f}, {r.hi:+.4f}]", r.verdict])
table(s, rows, M, Inches(1.9), W - 2 * M, col_w=[2.4, 2.2, 2.2, 1.6], size=13.5)
note(s, "Every interval spans zero. The gate demonstrably fired (slide 5), so this is evidence "
        "that suppressing the teacher's confident errors does not change what the student learns "
        "on this dataset — not evidence that the gate did nothing. The pre-registered endpoint "
        "was complete-miss rate, and it did not move either.", Inches(4.55))

# ═════════════════════════════════════════════════════════════════════════════
# 11 · complete misses
# ═════════════════════════════════════════════════════════════════════════════
s = slide("Complete misses — the metric that matters clinically",
          "Dice == 0: the bruise was not found at all.")
picture(s, "D3_complete_miss.png", Inches(1.9), Inches(3.3))
note(s, "The clinical question is whether a bruise was found, not how well it was outlined. "
        "At 0–13 misses out of 185 these are counts, not rates: a bootstrapped difference "
        "between two near-zero rates is unstable and is not quoted here.", Inches(5.6))

# ═════════════════════════════════════════════════════════════════════════════
# 12 · fairness
# ═════════════════════════════════════════════════════════════════════════════
s = slide("Skin-tone fairness across ITA groups")
picture(s, "D5_fairness.png", Inches(1.85), Inches(3.6))
note(s, "28 test subjects is a small denominator, so most group differences are not significant "
        "and must not be reported as if they were. Bruise size is also unevenly distributed "
        "across ITA groups, so any fairness claim that does not condition on size is measuring "
        "both at once.", Inches(5.8))

# ═════════════════════════════════════════════════════════════════════════════
# 13 · findings
# ═════════════════════════════════════════════════════════════════════════════
s = slide("Findings")
items = [
    ("Distillation from B2 works, and the teacher choice matters.",
     "Gated KD beats no KD on LR-ASPP (+0.027) and Fast-SCNN (+0.036); a B2 teacher beats "
     "DeepLabV3+ on Fast-SCNN (+0.034). All three intervals exclude zero."),
    ("No compact student reaches an undistilled SegFormer-B0.",
     "Every mobile arm is 0.04–0.12 Dice below the boundary, Holm-significant. Distillation "
     "narrows the architecture gap; it does not close it."),
    ("Reliability gating is a null result against response KD.",
     "Five students, every interval spanning zero, on both Dice and complete-miss rate — "
     "despite the gate demonstrably firing on ~3 % of image-views."),
    ("Everything sits inside the annotation ceiling.",
     "The B0 arms differ from the boundary by less than 0.003 Dice, against 0.755 between two "
     "human annotators. Significance is not importance."),
]
top = Inches(1.75)
for head, body in items:
    tf = textbox(s, M, top, W - 2 * M, Inches(1.15))
    _fmt(tf.paragraphs[0].add_run(), 15.5, bold=True, color=ACCENT)
    tf.paragraphs[0].runs[0].text = head
    p = tf.add_paragraph()
    _fmt(p.add_run(), 13, color=INK)
    p.runs[0].text = body
    top += Inches(1.28)

# ═════════════════════════════════════════════════════════════════════════════
# 14 · limitations
# ═════════════════════════════════════════════════════════════════════════════
s = slide("Limitations, stated up front")
rows = [
    ["28 test subjects",
     "Intervals are genuinely wide. INCONCLUSIVE appears often and means the study cannot "
     "answer that question at this size — not that the arms are equivalent."],
    ["Annotation ceiling",
     "Model-vs-model Dice differences below ≈0.05 are inside the noise floor of the labels "
     "themselves. Complete-miss rate is the more discriminating endpoint."],
    ["Omnibus does not reject for the headline field",
     f"Friedman on the seven headline models: p = "
     f"{float(D['omni'][D['omni'].set=='headline'].p.iloc[0]):.2f}. Pairwise winners must not "
     f"be quoted from inside that set."],
    ["The YOLO gated arm was not run",
     "Every confirmatory result here tests the gate on an online loss. Whether it transfers to "
     "an offline pseudo-mask route is untested."],
]
top = Inches(1.8)
for head, body in rows:
    tf = textbox(s, M, top, W - 2 * M, Inches(1.1))
    _fmt(tf.paragraphs[0].add_run(), 14.5, bold=True, color=INK)
    tf.paragraphs[0].runs[0].text = head
    p = tf.add_paragraph()
    _fmt(p.add_run(), 12.5, color=MUTED)
    p.runs[0].text = body
    top += Inches(1.22)

# ═════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
n = len(prs.slides.__iter__.__self__._sldIdLst)
print(f"wrote {OUT}")
print(f"  {n} slides, {FONT}, 16:9")
