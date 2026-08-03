#!/usr/bin/env python
"""Emit docs/b2_distillation_simple_median.pptx -- the median-Dice companion deck.

Same slides, same charts, same boundary logic as `72_generate_simple_b2_deck.py`.
The ONLY change is the endpoint: **median** per-image Dice everywhere instead of
mean.

WHY A SEPARATE SCRIPT RATHER THAN A FLAG
-----------------------------------------
The two decks are meant to be shown side by side, and the mean deck is already
built and checked. Parameterising 72 would put both at risk of one edit. This is a
build script, not a library -- clarity beats DRY here. If a third endpoint is ever
wanted, factor then.

WHAT CHANGES WITH THE ENDPOINT, AND WHY IT MATTERS
---------------------------------------------------
Median is NOT just "mean, but robust". Every number moves for a specific reason:

  * boundary       0.7663 -> 0.8129   (`segformer_b0_direct`)
  * human band     0.700-0.755 -> 0.750-0.809
  * axis           0.60-0.80 -> 0.65-0.85

The per-image Dice distribution is strongly left-skewed: a handful of complete
misses drag the mean down several points while the median barely moves. So median
Dice describes the TYPICAL image and is deliberately blind to the failure cases.

That is a feature for "how good is a normal prediction" and a liability for this
project, whose clinical endpoint is exactly the failure cases (handbook §1). Every
arm reads 0.032-0.066 Dice HIGHER on this deck than on the mean deck.

One thing that was checked rather than assumed: that uplift does **not** track
complete-miss count. Spearman(misses, median-mean gap) = -0.086, p = 0.84, and the
arm with the most misses (PP-MobileSeg response, 4) has the SMALLEST gap. With only
0-7 zeros out of 185 the misses are too few to move the mean much; the gap comes
from the broad left tail of low-but-nonzero scores. So the median is not hiding the
misses specifically -- it is hiding the whole weak tail. Miss counts are printed on
the mean-vs-median chart so the two endpoints can be read together.

Intervals are recomputed here with `significance.paired_contrast_multi`, which
bootstraps the median difference on the same resampled subject lists as the mean.
They are NOT read from `significance_contrast_family_stage_h.csv`, which stores
`delta_median` but drops `median_lo`/`median_hi`.
"""
from __future__ import annotations

from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib.patches import Patch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import sys
ROOT = Path(__file__).resolve().parent.parent
BUNDLE = ROOT / "BRUISE_UNIFIED"
RES = BUNDLE / "FINAL_RESULT"
FIG = RES / "figures" / "deck_median"
FIG.mkdir(parents=True, exist_ok=True)
OUT = ROOT / "docs" / "b2_distillation_simple_median.pptx"
sys.path.insert(0, str(BUNDLE))

from bruisekit import report as R          # noqa: E402
from bruisekit import significance as SG   # noqa: E402

plt.rcParams.update({
    "font.family": "Times New Roman", "font.size": 12,
    "axes.spines.top": False, "axes.spines.right": False,
    "axes.edgecolor": "#555555", "savefig.dpi": 220,
})

C_NONE, C_RESP, C_GATE, C_BAND = "#9AA3AD", "#1F3F6E", "#2E7D74", "#E8E4D9"
YLO, YHI = 0.65, 0.85

META = pd.read_csv(BUNDLE / "manifests" / "test.csv")
_cache: dict[str, pd.DataFrame] = {}


def per_image(m: str) -> pd.DataFrame:
    if m not in _cache:
        _cache[m] = R.normalize(pd.read_csv(RES / f"per_image_{m}.csv"), META)
    return _cache[m]


def med(m: str) -> float:
    return float(per_image(m).dice.median())


def misses(m: str) -> int:
    return int(per_image(m).complete_miss.sum())


REF = "segformer_b0_direct"
BOUND = med(REF)
ceil = pd.read_csv(RES / "annotation_ceiling.csv").set_index("comparison")
BAND = (float(ceil.loc["human: paul_vs_majority", "median_dice"]),
        float(ceil.loc["human: gbarimah_vs_erik", "median_dice"]))


def contrast(a: str, b: str) -> dict:
    """Median-Dice paired subject bootstrap, with a verdict on the median."""
    r = SG.paired_contrast_multi(per_image(a), per_image(b), a, b,
                                 n_boot=SG.N_BOOT_FINAL)
    return {"delta": r["delta_median"], "lo": r["median_lo"], "hi": r["median_hi"],
            "verdict": SG.verdict(r["delta_median"], r["median_lo"], r["median_hi"])}


SMALL = [("lraspp_mobilenetv3", "LR-ASPP\nMobileNetV3", "3.22 M"),
         ("topformer_tiny", "TopFormer-Tiny", "1.37 M"),
         ("ppmobileseg_tiny", "PP-MobileSeg-Tiny", "1.45 M")]


# ── charts ───────────────────────────────────────────────────────────────────
def _frame(ax):
    ax.axhspan(*BAND, color=C_BAND, zorder=0)
    ax.axhline(BOUND, color="#B03A2E", lw=1.7, ls=(0, (4, 3)), zorder=3)
    ax.set_ylim(YLO, YHI)
    ax.set_ylabel("median Dice", fontsize=12)
    ax.grid(axis="y", alpha=.20, ls=":")
    ax.set_axisbelow(True)


def _label(ax, xs, vals, size=11):
    for x, v in zip(xs, vals):
        ax.text(x, v + 0.004, f"{v:.3f}", ha="center", va="bottom", fontsize=size)


def _legend(ax, entries):
    h = [Patch(facecolor=c, label=l) for l, c in entries]
    h.append(plt.Line2D([], [], color="#B03A2E", lw=1.7, ls=(0, (4, 3)),
                        label=f"SegFormer-B0 direct, no teacher  ({BOUND:.3f})"))
    h.append(Patch(facecolor=C_BAND, label="human-vs-human agreement"))
    ax.legend(handles=h, frameon=False, fontsize=10.5, loc="upper center",
              bbox_to_anchor=(0.5, -0.13), ncol=2, handlelength=1.8,
              columnspacing=2.2)


def save(fig, name):
    fig.tight_layout()
    fig.savefig(FIG / f"{name}.png", bbox_inches="tight")
    plt.close(fig)
    print(f"  figures/deck_median/{name}.png")


def chart_b0():
    arms = [("segformer_b0_distilled", "Response KD\nfrom B2", C_RESP),
            ("segformer_b0_rgkd", "Reliability-gated KD\nfrom B2", C_GATE)]
    vals = [med(a) for a, _, _ in arms]
    fig, ax = plt.subplots(figsize=(8.6, 5.1))
    _frame(ax)
    xs = np.arange(len(arms))
    ax.bar(xs, vals, width=.42, color=[c for _, _, c in arms], zorder=2)
    _label(ax, xs, vals, 12.5)
    ax.set_xticks(xs); ax.set_xticklabels([l for _, l, _ in arms], fontsize=12)
    ax.set_xlim(-0.6, len(arms) - 0.4)
    _legend(ax, [("response KD", C_RESP), ("reliability-gated KD", C_GATE)])
    save(fig, "M1_b0_kd")


def chart_small(suffix, colour, label, name):
    base = [med(m) for m, _, _ in SMALL]
    arm = [med(f"{m}{suffix}") for m, _, _ in SMALL]
    xs, w = np.arange(len(SMALL)), .34
    fig, ax = plt.subplots(figsize=(9.6, 5.3))
    _frame(ax)
    ax.bar(xs - w / 2, base, width=w, color=C_NONE, zorder=2)
    ax.bar(xs + w / 2, arm, width=w, color=colour, zorder=2)
    _label(ax, xs - w / 2, base, 10.5); _label(ax, xs + w / 2, arm, 10.5)
    ax.set_xticks(xs)
    ax.set_xticklabels([f"{n}\n{p}" for _, n, p in SMALL], fontsize=11.5)
    ax.set_xlim(-0.6, len(SMALL) - 0.4)
    _legend(ax, [("no teacher", C_NONE), (label, colour)])
    save(fig, name)


def chart_head_to_head():
    models = [("segformer_b0", "SegFormer-B0", "segformer_b0_distilled", "segformer_b0_rgkd")] + \
             [(m, n, f"{m}_b2kd", f"{m}_rgkd") for m, n, _ in SMALL]
    xs, w = np.arange(len(models)), .34
    resp = [med(a) for _, _, a, _ in models]
    gate = [med(b) for _, _, _, b in models]
    fig, ax = plt.subplots(figsize=(10.4, 5.3))
    _frame(ax)
    ax.bar(xs - w / 2, resp, width=w, color=C_RESP, zorder=2)
    ax.bar(xs + w / 2, gate, width=w, color=C_GATE, zorder=2)
    _label(ax, xs - w / 2, resp, 10); _label(ax, xs + w / 2, gate, 10)
    ax.set_xticks(xs)
    ax.set_xticklabels([n for _, n, _, _ in models], fontsize=11.5)
    ax.set_xlim(-0.6, len(models) - 0.4)
    _legend(ax, [("response KD", C_RESP), ("reliability-gated KD", C_GATE)])
    save(fig, "M4_response_vs_gated")


def chart_mean_vs_median():
    """The one chart the mean deck does not have -- and the reason for the caveat.

    Mean and median side by side for every arm. The gap between the two bars is
    the complete-miss tax: it is what the median hides.
    """
    arms = [("segformer_b0_distilled", "B0\nresponse"), ("segformer_b0_rgkd", "B0\ngated"),
            ("lraspp_mobilenetv3_b2kd", "LR-ASPP\nresponse"), ("lraspp_mobilenetv3_rgkd", "LR-ASPP\ngated"),
            ("topformer_tiny_b2kd", "TopFormer\nresponse"), ("topformer_tiny_rgkd", "TopFormer\ngated"),
            ("ppmobileseg_tiny_b2kd", "PP-MobileSeg\nresponse"), ("ppmobileseg_tiny_rgkd", "PP-MobileSeg\ngated")]
    means = [float(per_image(a).dice.mean()) for a, _ in arms]
    meds = [med(a) for a, _ in arms]
    xs, w = np.arange(len(arms)), .36
    fig, ax = plt.subplots(figsize=(11.2, 5.3))
    ax.bar(xs - w / 2, means, width=w, color="#7E8C99", zorder=2)
    ax.bar(xs + w / 2, meds, width=w, color="#1F3F6E", zorder=2)
    for x, mn, md, (a, _) in zip(xs, means, meds, arms):
        ax.text(x, max(mn, md) + 0.006, f"+{md - mn:.3f}", ha="center", fontsize=9.5,
                color="#B03A2E")
        ax.text(x, YLO + 0.012, f"{misses(a)} miss", ha="center", fontsize=9,
                color="#555555")
    ax.set_ylim(YLO, YHI)
    ax.set_ylabel("Dice", fontsize=12)
    ax.set_xticks(xs); ax.set_xticklabels([n for _, n in arms], fontsize=10)
    ax.grid(axis="y", alpha=.20, ls=":"); ax.set_axisbelow(True)
    ax.legend(handles=[Patch(facecolor="#7E8C99", label="mean Dice"),
                       Patch(facecolor="#1F3F6E", label="median Dice")],
              frameon=False, fontsize=10.5, loc="upper center",
              bbox_to_anchor=(0.5, -0.13), ncol=2)
    save(fig, "M5_mean_vs_median")


def chart_fairness():
    fair = pd.read_csv(RES / "stage_h_fairness.csv").set_index("model")
    xs, w = np.arange(len(SMALL)), .26
    fig, ax = plt.subplots(figsize=(9.8, 5.3))
    for k, (lab, col, suf) in enumerate((("no teacher", C_NONE, ""),
                                         ("response KD", C_RESP, "_b2kd"),
                                         ("gated KD", C_GATE, "_rgkd"))):
        vals = [float(fair.loc[f"{m}{suf}", "fairness_gap"]) for m, _, _ in SMALL]
        sig = [bool(fair.loc[f"{m}{suf}", "significant"]) for m, _, _ in SMALL]
        pos = xs + (k - 1) * w
        ax.bar(pos, vals, width=w, color=col, zorder=2)
        for x, v, sg in zip(pos, vals, sig):
            ax.text(x, v + 0.004, f"{v:.3f}" + ("*" if sg else ""), ha="center",
                    fontsize=9.5, weight="bold" if sg else "normal")
    ax.set_xticks(xs); ax.set_xticklabels([n for _, n, _ in SMALL], fontsize=11.5)
    ax.set_ylabel("ITA fairness gap  (best group − worst group, median Dice)", fontsize=11)
    ax.set_ylim(0, 0.22); ax.grid(axis="y", alpha=.20, ls=":"); ax.set_axisbelow(True)
    ax.set_xlim(-0.6, len(SMALL) - 0.4)
    ax.legend(handles=[Patch(facecolor=c, label=l) for l, c in
                       (("no teacher", C_NONE), ("response KD", C_RESP), ("gated KD", C_GATE))],
              frameon=False, fontsize=10.5, loc="upper left")
    ax.text(0.99, 0.97, "* Kruskal–Wallis p < 0.05\nlower is fairer",
            transform=ax.transAxes, ha="right", va="top", fontsize=10, color="#555555")
    save(fig, "M6_fairness_small")


chart_b0()
chart_small("_b2kd", C_RESP, "response KD from B2", "M2_small_response")
chart_small("_rgkd", C_GATE, "reliability-gated KD from B2", "M3_small_gated")
chart_head_to_head()
chart_mean_vs_median()
chart_fairness()

# ═════════════════════════════════════════════════════════════════════════════
FONT = "Times New Roman"
INK, ACCENT, MUTED, RULE = (RGBColor(0x1A, 0x1A, 0x1A), RGBColor(0x1F, 0x3F, 0x6E),
                            RGBColor(0x6B, 0x6B, 0x6B), RGBColor(0xC9, 0xCF, 0xD8))
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H, M = prs.slide_width, prs.slide_height, Inches(0.75)
BLANK = prs.slide_layouts[6]


def _f(run, size, bold=False, color=INK, italic=False):
    run.font.name = FONT; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic; run.font.color.rgb = color


def tbox(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h); tb.text_frame.word_wrap = True
    return tb.text_frame


def slide(title, kicker=None):
    s = prs.slides.add_slide(BLANK)
    tf = tbox(s, M, Inches(0.40), W - 2 * M, Inches(0.7))
    _f(tf.paragraphs[0].add_run(), 25, True, ACCENT)
    tf.paragraphs[0].runs[0].text = title
    ln = s.shapes.add_shape(1, M, Inches(1.10), W - 2 * M, Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE
    ln.line.fill.background(); ln.shadow.inherit = False
    if kicker:
        kf = tbox(s, M, Inches(1.17), W - 2 * M, Inches(0.45))
        _f(kf.paragraphs[0].add_run(), 12.5, False, MUTED, True)
        kf.paragraphs[0].runs[0].text = kicker
    return s


def pic(s, name, top, max_h):
    from PIL import Image
    p = FIG / f"{name}.png"; iw, ih = Image.open(p).size
    w = W - 2 * M; h = int(w * ih / iw)
    if h > max_h:
        h, w = max_h, int(max_h * iw / ih)
    s.shapes.add_picture(str(p), int((W - w) / 2), top, width=int(w), height=int(h))


def note(s, text, top=Inches(6.5), size=11.5):
    tf = tbox(s, M, top, W - 2 * M, Inches(0.8))
    _f(tf.paragraphs[0].add_run(), size, False, MUTED)
    tf.paragraphs[0].runs[0].text = text


def table(s, rows, left, top, width, col_w, size=12, row_h=Inches(0.42)):
    t = s.shapes.add_table(len(rows), len(rows[0]), left, top, width,
                           row_h * len(rows)).table
    t.first_row = False
    tot = sum(col_w)
    for i, cw in enumerate(col_w):
        t.columns[i].width = int(width * cw / tot)
    for r, row in enumerate(rows):
        t.rows[r].height = row_h
        for c, v in enumerate(row):
            cell = t.cell(r, c); cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = 0
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run(); run.text = str(v)
            _f(run, size, r == 0, MUTED if r == 0 else INK)


# 1 title
s = prs.slides.add_slide(BLANK)
tf = tbox(s, M, Inches(2.6), W - 2 * M, Inches(1.2))
_f(tf.paragraphs[0].add_run(), 38, True, ACCENT)
tf.paragraphs[0].runs[0].text = "Distilling SegFormer-B2 into Compact Students"
tf2 = tbox(s, M, Inches(3.95), W - 2 * M, Inches(0.9))
_f(tf2.paragraphs[0].add_run(), 16)
tf2.paragraphs[0].runs[0].text = "Median-Dice view — the typical image, not the average"
ln = s.shapes.add_shape(1, M, Inches(5.05), Inches(2.2), Emu(19050))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT
ln.line.fill.background(); ln.shadow.inherit = False
tf3 = tbox(s, M, Inches(5.3), W - 2 * M, Inches(0.6))
_f(tf3.paragraphs[0].add_run(), 12.5, False, MUTED)
tf3.paragraphs[0].runs[0].text = "185 test images · 28 subjects · 3 seeds per arm · companion to the mean-Dice deck"

# 2 why median
s = slide("Why median, and what it hides",
          "Read alongside the mean deck, never instead of it.")
rows = [["", "mean Dice", "median Dice"],
        ["What it describes", "the average image, misses included", "the typical image"],
        ["Complete misses (Dice = 0)", "pull it down hard", "barely move it"],
        ["SegFormer-B0 direct (boundary)", "0.7663", f"{BOUND:.4f}"],
        ["Human-vs-human band", "0.700 – 0.755", f"{BAND[0]:.3f} – {BAND[1]:.3f}"]]
table(s, rows, M, Inches(1.8), W - 2 * M, [3.0, 3.0, 3.0], size=13, row_h=Inches(0.5))
note(s, "The per-image Dice distribution is strongly left-skewed, so a handful of complete "
        "misses drags the mean down several points while the median barely shifts. Median "
        "therefore answers “how good is a normal prediction” and is deliberately blind to the "
        "weak tail — which is where this project’s clinical endpoint lives. Every arm reads "
        "0.03–0.07 Dice higher here than on the mean deck. That is an artefact of the metric, "
        "not a result.", Inches(4.6))

# 3-5 charts
s = slide("SegFormer-B0: both KD methods sit on the boundary",
          "Dotted line = the same architecture trained with no teacher at all.")
pic(s, "M1_b0_kd", Inches(1.6), Inches(4.6))
c = contrast("segformer_b0_rgkd", "segformer_b0_distilled")
note(s, f"Gated − response on median Dice = {c['delta']:+.4f}, 95 % CI "
        f"[{c['lo']:+.4f}, {c['hi']:+.4f}] — spans zero, as it does on the mean.", Inches(6.35))

s = slide("Smaller students: response KD from B2",
          "Grey = no teacher. Navy = distilled from SegFormer-B2.")
pic(s, "M2_small_response", Inches(1.6), Inches(4.5))
note(s, "Each student improves on its own undistilled baseline. LR-ASPP reaches the human "
        "agreement band; none reaches the dotted boundary.", Inches(6.3))

s = slide("Smaller students: reliability-gated KD from B2",
          "Same teacher, same α, same recipe — only the loss differs.")
pic(s, "M3_small_gated", Inches(1.6), Inches(4.5))
note(s, "The gate fired on 2.8 % of image-views. On median Dice, as on mean, the bars land in "
        "essentially the same place as response KD.", Inches(6.3))

# 6 head to head + table
s = slide("Response KD vs reliability-gated KD, head to head")
pic(s, "M4_response_vs_gated", Inches(1.5), Inches(4.3))
note(s, "Miss counts are printed on the next slide — they are what this chart cannot show.",
     Inches(6.15))

s = slide("The same comparison as numbers",
          "Median-Dice difference, paired subject-level bootstrap, 28 subjects, 10 000 draws.")
rows = [["Student", "Δ median (gated − response)", "95 % CI", "Verdict"]]
for m, n in [("segformer_b0", "SegFormer-B0"), ("lraspp_mobilenetv3", "LR-ASPP MobileNetV3"),
             ("topformer_tiny", "TopFormer-Tiny"), ("ppmobileseg_tiny", "PP-MobileSeg-Tiny")]:
    a = f"{m}_rgkd"
    b = "segformer_b0_distilled" if m == "segformer_b0" else f"{m}_b2kd"
    c = contrast(a, b)
    rows.append([n, f"{c['delta']:+.4f}", f"[{c['lo']:+.4f}, {c['hi']:+.4f}]", c["verdict"]])
table(s, rows, M, Inches(1.8), W - 2 * M, [2.6, 2.4, 2.2, 1.6], size=13, row_h=Inches(0.45))

rows2 = [["Against the boundary (SegFormer-B0 direct)", "Δ median", "95 % CI", "Verdict"]]
for arm, n in [("segformer_b0_distilled", "SegFormer-B0 · response"),
               ("segformer_b0_rgkd", "SegFormer-B0 · gated"),
               ("lraspp_mobilenetv3_b2kd", "LR-ASPP · response"),
               ("topformer_tiny_b2kd", "TopFormer · response"),
               ("ppmobileseg_tiny_b2kd", "PP-MobileSeg · response")]:
    c = contrast(arm, REF)
    rows2.append([n, f"{c['delta']:+.4f}", f"[{c['lo']:+.4f}, {c['hi']:+.4f}]", c["verdict"]])
table(s, rows2, M, Inches(4.15), W - 2 * M, [2.6, 2.4, 2.2, 1.6], size=13, row_h=Inches(0.45))

# 7 mean vs median
s = slide("What the median hides",
          "Red = how much the median flatters each arm. Miss counts along the bottom.")
pic(s, "M5_mean_vs_median", Inches(1.6), Inches(4.5))
note(s, "The gap between the two bars is the weak tail the median discards — 0.032 to 0.066 "
        "for every arm. It does NOT track complete misses (Spearman −0.09, p = 0.84): with only "
        "0–7 zeros out of 185 the misses are too few to move the mean much, so what the median "
        "hides is the broad tail of low-but-nonzero scores. Quote both endpoints, never one.",
     Inches(6.25))

# 8 fairness
s = slide("Skin-tone fairness across ITA groups",
          "Gap between best and worst ITA group, on median Dice. Lower is fairer.")
pic(s, "M6_fairness_small", Inches(1.6), Inches(4.5))
note(s, "This chart is unchanged from the mean deck — the fairness gap is defined on median "
        "Dice in both. PP-MobileSeg narrows under KD and loses significance when gated; "
        "TopFormer moves the wrong way; LR-ASPP is never significant.", Inches(6.3))

# 9 findings
s = slide("What the results say")
for head, body, top in [
    ("Distillation from B2 helps every compact student.",
     "Each smaller model beats its own undistilled baseline on median Dice, and LR-ASPP "
     "reaches the human-agreement band.", Inches(1.7)),
    ("None of them reaches an undistilled SegFormer-B0.",
     "Gaps of 0.035–0.105 median Dice to the dotted boundary; TopFormer and PP-MobileSeg "
     "are INFERIOR with intervals excluding the margin.", Inches(2.96)),
    ("Reliability gating did not beat response KD.",
     "Same conclusion as the mean deck — every interval spans zero. The endpoint does not "
     "change the answer, which is itself reassuring.", Inches(4.22)),
    ("Median flatters every arm by roughly the same amount.",
     "0.032–0.066 Dice higher than the mean, and the uplift does not track complete misses "
     "(Spearman −0.09). The ranking is unchanged; the level is not.", Inches(5.48))]:
    tf = tbox(s, M, top, W - 2 * M, Inches(1.15))
    _f(tf.paragraphs[0].add_run(), 15, True, ACCENT)
    tf.paragraphs[0].runs[0].text = head
    p = tf.add_paragraph(); _f(p.add_run(), 12.5)
    p.runs[0].text = body

prs.save(OUT)
print(f"\nwrote {OUT}")
print(f"  {len(prs.slides._sldIdLst)} slides, {FONT}, 16:9")
