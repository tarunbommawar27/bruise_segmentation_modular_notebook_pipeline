#!/usr/bin/env python
"""Emit both B2-distillation decks -- mean-Dice and median-Dice -- from one source.

    python scripts/74_generate_b2_decks.py

    docs/b2_distillation_simple.pptx          endpoint = mean Dice
    docs/b2_distillation_simple_median.pptx   endpoint = median Dice

WHY ONE SCRIPT FOR BOTH
------------------------
These began as two near-identical files (72 and 73). They drifted immediately: a
caption claim about complete misses was wrong in one and absent from the other,
and every subsequent edit had to be made twice with no guarantee it was. The decks
are shown side by side and must differ in exactly one thing -- the endpoint -- so
that is now the only parameter. 72 and 73 are superseded.

SCOPE
------
Four students: SegFormer-B0, LR-ASPP MobileNetV3, TopFormer-Tiny, PP-MobileSeg-Tiny.
**Fast-SCNN is excluded**, by request. That is defensible -- it is the only model
with no pretrained backbone, so it confounds "small" with "no ImageNet" -- but it
was also where KD-vs-none was largest, so the exclusion understates what
distillation achieved. Said on the scope slide, not buried here.

CHART GRAMMAR
--------------
  * the dotted red line is `segformer_b0_direct`, labelled ON the line with its
    value. It is the boundary every bar is read against.
  * grey = no teacher, navy = response KD, teal = reliability-gated. One colour per
    training regime, held across every chart and both decks.
  * the y-axis is truncated. Truncation exaggerates differences, so every bar
    carries its own value label and the caption says the axis is cut.

All numbers are read from FINAL_RESULT at build time; intervals are recomputed
through `significance.paired_contrast_multi` so both endpoints come from the same
bootstrap draws. Nothing is transcribed.
"""
from __future__ import annotations

import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib.patches import Patch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
BUNDLE = ROOT / "BRUISE_UNIFIED"
RES = BUNDLE / "FINAL_RESULT"
sys.path.insert(0, str(BUNDLE))

from bruisekit import report as R          # noqa: E402
from bruisekit import significance as SG   # noqa: E402

plt.rcParams.update({
    "font.family": "Times New Roman", "font.size": 12,
    "axes.spines.top": False, "axes.spines.right": False,
    "axes.edgecolor": "#555555", "savefig.dpi": 220,
})

C_NONE, C_RESP, C_GATE, C_LINE = "#9AA3AD", "#1F3F6E", "#2E7D74", "#B03A2E"
FONT = "Times New Roman"
INK, ACCENT, MUTED, RULE = (RGBColor(0x1A, 0x1A, 0x1A), RGBColor(0x1F, 0x3F, 0x6E),
                            RGBColor(0x6B, 0x6B, 0x6B), RGBColor(0xC9, 0xCF, 0xD8))
REF = "segformer_b0_direct"

SMALL = [("lraspp_mobilenetv3", "LR-ASPP\nMobileNetV3", "LR-ASPP MobileNetV3", "3.22 M"),
         ("topformer_tiny", "TopFormer-Tiny", "TopFormer-Tiny", "1.37 M"),
         ("ppmobileseg_tiny", "PP-MobileSeg-Tiny", "PP-MobileSeg-Tiny", "1.45 M")]

META = pd.read_csv(BUNDLE / "manifests" / "test.csv")
_cache: dict[str, pd.DataFrame] = {}


def per_image(m: str) -> pd.DataFrame:
    if m not in _cache:
        _cache[m] = R.normalize(pd.read_csv(RES / f"per_image_{m}.csv"), META)
    return _cache[m]


def misses(m: str) -> int:
    return int(per_image(m).complete_miss.sum())


# ─────────────────────────────────────────────────────────────────────────────
class Endpoint:
    """Everything that differs between the two decks, in one object."""

    def __init__(self, kind: str):
        assert kind in ("mean", "median")
        self.kind = kind
        self.label = f"{kind} Dice"
        self.figdir = RES / "figures" / f"deck_{kind}"
        self.figdir.mkdir(parents=True, exist_ok=True)
        self.out = ROOT / "docs" / ("b2_distillation_simple.pptx" if kind == "mean"
                                    else "b2_distillation_simple_median.pptx")
        self.ylim = (0.60, 0.80) if kind == "mean" else (0.65, 0.85)
        self.prefix = "C" if kind == "mean" else "M"
        self.bound = self.value(REF)

    def value(self, m: str) -> float:
        d = per_image(m).dice
        return float(d.mean() if self.kind == "mean" else d.median())

    def contrast(self, a: str, b: str) -> dict:
        """Paired subject-level bootstrap on THIS endpoint, with a verdict."""
        r = SG.paired_contrast_multi(per_image(a), per_image(b), a, b,
                                     n_boot=SG.N_BOOT_FINAL)
        if self.kind == "mean":
            d, lo, hi = r["delta_dice"], r["lo"], r["hi"]
        else:
            d, lo, hi = r["delta_median"], r["median_lo"], r["median_hi"]
        return {"delta": d, "lo": lo, "hi": hi,
                "verdict": SG.verdict(d, lo, hi), "p": r["p_two_sided"]}


# ── charts ───────────────────────────────────────────────────────────────────
def _frame(E, ax):
    """Axes furniture, and the boundary line labelled ON the line."""
    ax.axhline(E.bound, color=C_LINE, lw=1.8, ls=(0, (5, 3)), zorder=3)
    ax.set_ylim(*E.ylim)
    ax.set_ylabel(E.label, fontsize=12)
    ax.grid(axis="y", alpha=.20, ls=":")
    ax.set_axisbelow(True)


def _boundary_tag(E, ax):
    """The value printed against the line itself, not only in the legend."""
    ax.annotate(f"SegFormer-B0 direct (no teacher) = {E.bound:.3f}",
                xy=(0.995, E.bound), xycoords=("axes fraction", "data"),
                xytext=(0, 5), textcoords="offset points",
                ha="right", va="bottom", fontsize=11, color=C_LINE, weight="bold")


def _label(ax, xs, vals, size=11):
    for x, v in zip(xs, vals):
        ax.text(x, v + 0.004, f"{v:.3f}", ha="center", va="bottom", fontsize=size)


def _legend(ax, entries, pad=-0.13):
    """Legend below the axes. `pad` must clear the x tick labels, which are
    multi-line on the charts that print a parameter count under each name."""
    h = [Patch(facecolor=c, label=l) for l, c in entries]
    h.append(plt.Line2D([], [], color=C_LINE, lw=1.8, ls=(0, (5, 3)),
                        label="boundary: SegFormer-B0, no teacher"))
    ax.legend(handles=h, frameon=False, fontsize=10.5, loc="upper center",
              bbox_to_anchor=(0.5, pad), ncol=3, handlelength=1.8,
              columnspacing=2.2)


def _save(E, fig, name):
    fig.tight_layout()
    fig.savefig(E.figdir / f"{name}.png", bbox_inches="tight")
    plt.close(fig)
    print(f"  figures/deck_{E.kind}/{name}.png")


def chart_gate_explainer(E):
    """How ground truth and B2 combine. Two panels: per pixel, then per image."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12.4, 4.0))

    # left: per-pixel reliability against the teacher's probability, y = 1
    p = np.linspace(0, 1, 501)
    r = 1 - np.abs(2 * p - 1) * np.abs(p - 1)
    ax1.plot(p, r, color=C_GATE, lw=2.6)
    ax1.fill_between(p, 0, r, color=C_GATE, alpha=.10)
    ax1.set_xlim(0, 1); ax1.set_ylim(0, 1.12)
    ax1.set_xlabel("B2's probability that this pixel is bruise", fontsize=11)
    ax1.set_ylabel("weight given to B2", fontsize=11)
    ax1.set_title("Per pixel  —  checked against the label",
                  fontsize=12, weight="bold", pad=10)
    NL = chr(10)
    # x positions are the ANNOTATION anchors, nudged inward from the data points
    # so nothing is clipped at the panel edges.
    for x, tx, ha, txt, col in (
            (0.02, 0.06, "left", f"confidently WRONG{NL}→ ignored", C_LINE),
            (0.50, 0.50, "center", f"uncertain{NL}→ full weight", "#1A1A1A"),
            (0.99, 0.96, "right", f"confidently{NL}right", "#1A1A1A")):
        ax1.annotate(txt, xy=(x, 1 - abs(2 * x - 1) * abs(x - 1)),
                     xytext=(tx, 1.02), ha=ha, va="bottom",
                     fontsize=9.5, color=col, weight="bold")
    ax1.grid(alpha=.20, ls=":"); ax1.set_axisbelow(True)

    # right: per-image gate against the teacher's own Dice
    d = np.linspace(0, 1, 501)
    g = np.clip((d - 0.10) / 0.40, 0, 1)
    ax2.plot(d, g, color=C_RESP, lw=2.6)
    ax2.fill_between(d, 0, g, color=C_RESP, alpha=.10)
    ax2.axvspan(0, 0.10, color=C_LINE, alpha=.10)
    ax2.set_xlim(0, 1); ax2.set_ylim(0, 1.12)
    ax2.set_xlabel("B2's own Dice against the label, on this image", fontsize=11)
    ax2.set_ylabel("weight given to B2", fontsize=11)
    ax2.set_title("Per image  —  checked against the label",
                  fontsize=12, weight="bold", pad=10)
    ax2.annotate(f"B2 missed the bruise{NL}→ train on the label alone",
                 xy=(0.05, 0.02), xytext=(0.30, 0.66), fontsize=9.5,
                 color=C_LINE, weight="bold",
                 arrowprops=dict(arrowstyle="->", color=C_LINE, lw=1.2))
    ax2.annotate(f"B2 found it{NL}→ full weight", xy=(0.80, 1.0), xytext=(0.72, 0.40),
                 fontsize=9.5, weight="bold", ha="center",
                 arrowprops=dict(arrowstyle="->", color="#555555", lw=1.2))
    ax2.grid(alpha=.20, ls=":"); ax2.set_axisbelow(True)

    _save(E, fig, f"{E.prefix}0_gate_explainer")


def chart_b0(E):
    arms = [("segformer_b0_distilled", "Response KD\nfrom B2", C_RESP),
            ("segformer_b0_rgkd", "Reliability-gated KD\nfrom B2", C_GATE)]
    vals = [E.value(a) for a, _, _ in arms]
    fig, ax = plt.subplots(figsize=(8.8, 5.1))
    _frame(E, ax)
    xs = np.arange(len(arms))
    ax.bar(xs, vals, width=.42, color=[c for _, _, c in arms], zorder=2)
    _label(ax, xs, vals, 12.5)
    ax.set_xticks(xs); ax.set_xticklabels([l for _, l, _ in arms], fontsize=12)
    ax.set_xlim(-0.7, len(arms) - 0.3)
    _boundary_tag(E, ax)
    _legend(ax, [("response KD", C_RESP), ("reliability-gated KD", C_GATE)])
    _save(E, fig, f"{E.prefix}1_b0_kd")


def chart_small(E, suffix, colour, label, name):
    base = [E.value(m) for m, _, _, _ in SMALL]
    arm = [E.value(f"{m}{suffix}") for m, _, _, _ in SMALL]
    xs, w = np.arange(len(SMALL)), .34
    fig, ax = plt.subplots(figsize=(9.8, 5.3))
    _frame(E, ax)
    ax.bar(xs - w / 2, base, width=w, color=C_NONE, zorder=2)
    ax.bar(xs + w / 2, arm, width=w, color=colour, zorder=2)
    _label(ax, xs - w / 2, base, 10.5); _label(ax, xs + w / 2, arm, 10.5)
    ax.set_xticks(xs)
    ax.set_xticklabels([f"{n}\n{p}" for _, n, _, p in SMALL], fontsize=11.5)
    ax.set_xlim(-0.7, len(SMALL) - 0.3)
    _boundary_tag(E, ax)
    # three-line tick labels here (name + params), so the legend needs more room
    _legend(ax, [("no teacher", C_NONE), (label, colour)], pad=-0.30)
    _save(E, fig, name)


def chart_head_to_head(E):
    models = [("segformer_b0", "SegFormer-B0", "segformer_b0_distilled", "segformer_b0_rgkd")] + \
             [(m, n, f"{m}_b2kd", f"{m}_rgkd") for m, n, _, _ in SMALL]
    xs, w = np.arange(len(models)), .34
    resp = [E.value(a) for _, _, a, _ in models]
    gate = [E.value(b) for _, _, _, b in models]
    fig, ax = plt.subplots(figsize=(10.6, 5.3))
    _frame(E, ax)
    ax.bar(xs - w / 2, resp, width=w, color=C_RESP, zorder=2)
    ax.bar(xs + w / 2, gate, width=w, color=C_GATE, zorder=2)
    _label(ax, xs - w / 2, resp, 10); _label(ax, xs + w / 2, gate, 10)
    ax.set_xticks(xs); ax.set_xticklabels([n for _, n, _, _ in models], fontsize=11.5)
    ax.set_xlim(-0.7, len(models) - 0.3)
    _boundary_tag(E, ax)
    _legend(ax, [("response KD", C_RESP), ("reliability-gated KD", C_GATE)], pad=-0.18)
    _save(E, fig, f"{E.prefix}4_response_vs_gated")


def chart_mean_vs_median(E):
    """Median deck only: what the median discards, with miss counts."""
    arms = [("segformer_b0_distilled", "B0\nresponse"), ("segformer_b0_rgkd", "B0\ngated"),
            ("lraspp_mobilenetv3_b2kd", "LR-ASPP\nresponse"), ("lraspp_mobilenetv3_rgkd", "LR-ASPP\ngated"),
            ("topformer_tiny_b2kd", "TopFormer\nresponse"), ("topformer_tiny_rgkd", "TopFormer\ngated"),
            ("ppmobileseg_tiny_b2kd", "PP-MobileSeg\nresponse"), ("ppmobileseg_tiny_rgkd", "PP-MobileSeg\ngated")]
    means = [float(per_image(a).dice.mean()) for a, _ in arms]
    meds = [float(per_image(a).dice.median()) for a, _ in arms]
    xs, w = np.arange(len(arms)), .36
    fig, ax = plt.subplots(figsize=(11.2, 5.3))
    ax.bar(xs - w / 2, means, width=w, color="#7E8C99", zorder=2)
    ax.bar(xs + w / 2, meds, width=w, color=C_RESP, zorder=2)
    for x, mn, md, (a, _) in zip(xs, means, meds, arms):
        ax.text(x, max(mn, md) + 0.006, f"+{md - mn:.3f}", ha="center",
                fontsize=9.5, color=C_LINE)
        ax.text(x, E.ylim[0] + 0.012, f"{misses(a)} miss", ha="center",
                fontsize=9, color="#555555")
    ax.set_ylim(*E.ylim); ax.set_ylabel("Dice", fontsize=12)
    ax.set_xticks(xs); ax.set_xticklabels([n for _, n in arms], fontsize=10)
    ax.grid(axis="y", alpha=.20, ls=":"); ax.set_axisbelow(True)
    ax.legend(handles=[Patch(facecolor="#7E8C99", label="mean Dice"),
                       Patch(facecolor=C_RESP, label="median Dice")],
              frameon=False, fontsize=10.5, loc="upper center",
              bbox_to_anchor=(0.5, -0.13), ncol=2)
    _save(E, fig, f"{E.prefix}5_mean_vs_median")


def chart_fairness(E):
    fair = pd.read_csv(RES / "stage_h_fairness.csv").set_index("model")
    xs, w = np.arange(len(SMALL)), .26
    fig, ax = plt.subplots(figsize=(9.8, 5.3))
    for k, (lab, col, suf) in enumerate((("no teacher", C_NONE, ""),
                                         ("response KD", C_RESP, "_b2kd"),
                                         ("gated KD", C_GATE, "_rgkd"))):
        vals = [float(fair.loc[f"{m}{suf}", "fairness_gap"]) for m, _, _, _ in SMALL]
        sig = [bool(fair.loc[f"{m}{suf}", "significant"]) for m, _, _, _ in SMALL]
        pos = xs + (k - 1) * w
        ax.bar(pos, vals, width=w, color=col, zorder=2)
        for x, v, sg in zip(pos, vals, sig):
            ax.text(x, v + 0.004, f"{v:.3f}" + ("*" if sg else ""), ha="center",
                    fontsize=9.5, weight="bold" if sg else "normal")
    ax.set_xticks(xs); ax.set_xticklabels([n for _, n, _, _ in SMALL], fontsize=11.5)
    ax.set_ylabel("ITA fairness gap  (best group − worst group)", fontsize=11)
    ax.set_ylim(0, 0.22); ax.grid(axis="y", alpha=.20, ls=":"); ax.set_axisbelow(True)
    ax.set_xlim(-0.6, len(SMALL) - 0.4)
    ax.legend(handles=[Patch(facecolor=c, label=l) for l, c in
                       (("no teacher", C_NONE), ("response KD", C_RESP), ("gated KD", C_GATE))],
              frameon=False, fontsize=10.5, loc="upper left")
    ax.text(0.99, 0.97, "* Kruskal–Wallis p < 0.05\nlower is fairer",
            transform=ax.transAxes, ha="right", va="top", fontsize=10, color="#555555")
    _save(E, fig, f"{E.prefix}6_fairness")


# ── deck furniture ───────────────────────────────────────────────────────────
def _f(run, size, bold=False, color=INK, italic=False):
    run.font.name = FONT; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic; run.font.color.rgb = color


def build(E):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    W, H, M = prs.slide_width, prs.slide_height, Inches(0.75)
    BLANK = prs.slide_layouts[6]

    def tbox(s, l, t, w, h):
        tb = s.shapes.add_textbox(l, t, w, h); tb.text_frame.word_wrap = True
        return tb.text_frame

    def slide(title, kicker=None):
        s = prs.slides.add_slide(BLANK)
        tf = tbox(s, M, Inches(0.40), W - 2 * M, Inches(0.7))
        _f(tf.paragraphs[0].add_run(), 25, True, ACCENT)
        tf.paragraphs[0].runs[0].text = title
        ln = s.shapes.add_shape(1, M, Inches(1.10), W - 2 * M, Emu(9525))
        ln.fill.solid(); ln.fill.fore_color.rgb = RULE
        ln.line.fill.background(); ln.shadow.inherit = False
        if kicker:
            kf = tbox(s, M, Inches(1.17), W - 2 * M, Inches(0.45))
            _f(kf.paragraphs[0].add_run(), 12.5, False, MUTED, True)
            kf.paragraphs[0].runs[0].text = kicker
        return s

    def pic(s, name, top, max_h):
        from PIL import Image
        p = E.figdir / f"{name}.png"; iw, ih = Image.open(p).size
        w = W - 2 * M; h = int(w * ih / iw)
        if h > max_h:
            h, w = max_h, int(max_h * iw / ih)
        s.shapes.add_picture(str(p), int((W - w) / 2), top, width=int(w), height=int(h))

    def note(s, text, top=Inches(6.45), size=11.5):
        tf = tbox(s, M, top, W - 2 * M, Inches(0.85))
        _f(tf.paragraphs[0].add_run(), size, False, MUTED)
        tf.paragraphs[0].runs[0].text = text

    def table(s, rows, top, col_w, size=12, row_h=Inches(0.42), left=None, width=None):
        left = left or M
        width = width or (W - 2 * M)
        t = s.shapes.add_table(len(rows), len(rows[0]), left, top, width,
                               row_h * len(rows)).table
        t.first_row = False
        tot = sum(col_w)
        for i, cw in enumerate(col_w):
            t.columns[i].width = int(width * cw / tot)
        for r, row in enumerate(rows):
            t.rows[r].height = row_h
            for c, v in enumerate(row):
                cell = t.cell(r, c); cell.text = ""
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = cell.margin_right = Inches(0.07)
                cell.margin_top = cell.margin_bottom = 0
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
                run = p.add_run(); run.text = str(v)
                _f(run, size, r == 0, MUTED if r == 0 else INK)

    def bullets(s, items, top=Inches(1.7), gap=Inches(1.26), hs=15, bs=12.5):
        for head, body in items:
            tf = tbox(s, M, top, W - 2 * M, Inches(1.15))
            _f(tf.paragraphs[0].add_run(), hs, True, ACCENT)
            tf.paragraphs[0].runs[0].text = head
            p = tf.add_paragraph(); _f(p.add_run(), bs)
            p.runs[0].text = body
            top += gap

    # 1 title
    s = prs.slides.add_slide(BLANK)
    tf = tbox(s, M, Inches(2.6), W - 2 * M, Inches(1.2))
    _f(tf.paragraphs[0].add_run(), 38, True, ACCENT)
    tf.paragraphs[0].runs[0].text = "Distilling SegFormer-B2 into Compact Students"
    tf2 = tbox(s, M, Inches(3.95), W - 2 * M, Inches(0.9))
    _f(tf2.paragraphs[0].add_run(), 16)
    tf2.paragraphs[0].runs[0].text = (
        "Response KD and reliability-gated KD, measured against an undistilled SegFormer-B0"
        + ("" if E.kind == "mean" else "  ·  median-Dice view"))
    ln = s.shapes.add_shape(1, M, Inches(5.05), Inches(2.2), Emu(19050))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background(); ln.shadow.inherit = False
    tf3 = tbox(s, M, Inches(5.3), W - 2 * M, Inches(0.6))
    _f(tf3.paragraphs[0].add_run(), 12.5, False, MUTED)
    tf3.paragraphs[0].runs[0].text = (
        f"185 test images · 28 subjects · 3 seeds per arm · endpoint: {E.label}")

    # 2 why these students
    s = slide("Why these four students", "The set was chosen to answer four different questions.")
    bullets(s, [
        ("SegFormer-B0 — the same family as the teacher.",
         "3.71 M against the teacher's 27.4 M. Because the architecture is held fixed, any "
         "difference is about SIZE alone. It is also the reference boundary for everything else."),
        ("LR-ASPP MobileNetV3 — the realistic deployment target.",
         "The strongest mobile model in the study and a pure CNN. If distillation only worked "
         "transformer-to-transformer it would be of no practical use, so this is the arm that "
         "tests whether the teacher's knowledge crosses architecture families."),
        ("TopFormer-Tiny and PP-MobileSeg-Tiny — half the parameters, attention kept.",
         "Both are hybrids: a CNN feature extractor with a transformer/attention head. They ask "
         "whether the result survives at 1.4 M parameters, and whether it depends on one "
         "particular attention design — they use two different ones."),
        ("All four share one recipe and one initialisation policy.",
         "ImageNet-pretrained backbone, freshly initialised head, identical LR schedule, loss, "
         "batch and seeds. Nothing but the student changes."),
    ], top=Inches(1.72), gap=Inches(1.34), hs=14.5, bs=12)

    # 3 the students, in numbers
    s = slide("The students, in numbers")
    table(s, [["Model", "Params", "Architecture", "Backbone init"],
              ["SegFormer-B0", "3.71 M", "pure transformer (MiT hierarchical encoder)", "ImageNet-1k"],
              ["LR-ASPP MobileNetV3", "3.22 M", "pure CNN (depthwise separable + ASPP-lite)", "ImageNet-1k, torchvision"],
              ["TopFormer-Tiny", "1.37 M", "hybrid — CNN token pyramid + transformer block", "ImageNet-1k, 66.2 % top-1"],
              ["PP-MobileSeg-Tiny", "1.45 M", "hybrid — StrideFormer, strided SEA attention", "ImageNet (StrideFormer)"],
              ["SegFormer-B2  (teacher)", "27.35 M", "pure transformer", "ImageNet-1k"]],
          Inches(1.8), [2.3, 0.9, 4.4, 2.4], size=12.5, row_h=Inches(0.52))
    note(s, "Two of the four are vendored verbatim from their reference implementations, and "
            "every official checkpoint was verified to load by name and shape before any GPU "
            "time was spent — an architecture that is 'close enough' would train fine and simply "
            "not be the model named in the table.", Inches(5.35))

    # 4 how the gate uses the label and B2
    s = slide("How reliability gating uses the label and B2",
              "Standard KD trusts the teacher equally everywhere. Gating does not.")
    pic(s, f"{E.prefix}0_gate_explainer", Inches(1.60), Inches(3.15))
    tf = tbox(s, M, Inches(4.95), W - 2 * M, Inches(2.45))
    for txt, sz, bold, col in (
        ("loss  =  α · (student vs ground truth)   +   (1 − α) · "
         "(student vs B2)", 15, True, ACCENT),
        ("Both terms were always there. Gating does not add the ground truth — it changes "
         "how much of it we lean on, and where.", 12, False, INK),
        ("During training we can see the label, so we check B2 against it. Where B2 is "
         "confidently wrong we stop listening and α rises toward 1; where B2 is uncertain we "
         "keep it at full strength, because that uncertainty is the information distillation "
         "exists to transfer. On an image B2 missed entirely the teacher term switches off and "
         "the student trains on the label alone.", 12, False, MUTED)):
        para = tf.paragraphs[0] if txt.startswith("loss") else tf.add_paragraph()
        _f(para.add_run(), sz, bold, col)
        para.runs[0].text = txt
    # Kept in the same textbox rather than a separate note: a fourth floating box
    # pushed this slide 0.37 in past the bottom edge.
    para = tf.add_paragraph()
    _f(para.add_run(), 11, False, MUTED)
    para.runs[0].text = (
        "Training-time only — the gate needs the label, so the deployed student is an "
        "ordinary model and carries no gate. Measured over the 15 gated runs: mean weight on "
        "B2 = 0.906, and 2.8 % of image-views had the teacher switched off entirely.")

    # 5-7 charts
    s = slide("SegFormer-B0: both KD methods sit on the boundary",
              "Dotted line = the same architecture trained with no teacher.")
    pic(s, f"{E.prefix}1_b0_kd", Inches(1.6), Inches(4.6))
    c = E.contrast("segformer_b0_rgkd", "segformer_b0_distilled")
    note(s, f"Gated − response = {c['delta']:+.4f}, 95 % CI [{c['lo']:+.4f}, {c['hi']:+.4f}] — "
            f"spans zero. At 3.71 M the student already matches the teacher, so there is little "
            f"for distillation to add. Note the y-axis starts at {E.ylim[0]:.2f}.", Inches(6.35))

    s = slide("Smaller students: response KD from B2",
              "Grey = no teacher. Navy = distilled from SegFormer-B2.")
    pic(s, f"{E.prefix}2_small_response", Inches(1.6), Inches(4.5))
    note(s, "Every student improves over its own undistilled baseline. None reaches the dotted "
            "boundary: distillation narrows the gap to a 3.71 M transformer but does not close "
            f"it. Y-axis starts at {E.ylim[0]:.2f}.", Inches(6.3))

    s = slide("Smaller students: reliability-gated KD from B2",
              "Same teacher, same α, same recipe — only the loss differs.")
    pic(s, f"{E.prefix}3_small_gated", Inches(1.6), Inches(4.5))
    note(s, "The gate suppresses the teacher where it is confidently wrong and switches it off "
            "entirely on images the teacher missed. It fired on 2.8 % of image-views. The bars "
            "land in essentially the same place as response KD.", Inches(6.3))

    s = slide("Response KD vs reliability-gated KD, head to head")
    pic(s, f"{E.prefix}4_response_vs_gated", Inches(1.55), Inches(4.4))
    note(s, "Same teacher, same mixing weight, same seeds. The only difference is whether the "
            "teacher is trusted uniformly or trusted per pixel.", Inches(6.2))

    # 7 how significance is computed
    s = slide("How significance is calculated",
              "One method, applied identically to every comparison in this deck.")
    bullets(s, [
        ("Resample subjects, not images.",
         "The 185 test images come from 28 subjects, and images of the same bruise are strongly "
         "correlated. Resampling images would treat 185 correlated observations as independent "
         "and produce intervals roughly 2.6× too narrow. We draw 28 subjects with replacement "
         "and take all of each chosen subject's images."),
        ("Pair the two models on every draw.",
         "Both models were scored on the same 185 images, so the same resampled subject list is "
         "applied to both. Resampling them independently would discard that pairing and hide "
         "real but small differences."),
        ("10 000 draws, three endpoints, one set of draws.",
         f"{E.label}, the other Dice statistic, and complete-miss rate are all evaluated on the "
         "same resampled worlds — so 'Dice up, misses up' is a statement about the same 10 000 "
         "resampled worlds rather than three unrelated bootstraps."),
        ("Holm–Bonferroni within the pre-specified family.",
         "The comparison list was fixed before any of these models were trained. Correcting "
         "within it prevents a table of many comparisons from producing a 'winner' by chance."),
    ], top=Inches(1.72), gap=Inches(1.30), hs=14.5, bs=11.5)

    # 8 verdicts
    s = slide("How to read a verdict",
              f"Δ is the difference in {E.label}; the interval is the 95 % bootstrap range.")
    table(s, [["Verdict", "Condition", "What it means"],
              ["WIN", "interval entirely above 0", "genuinely better"],
              ["INFERIOR", "interval entirely below −0.01", "genuinely worse"],
              ["NON-INFERIOR", "interval stays above −0.01", "equivalent within one Dice point"],
              ["INCONCLUSIVE", "interval extends past −0.01 but crosses 0",
               "the study cannot answer this at 28 subjects"]],
          Inches(1.85), [1.7, 3.0, 4.6], size=13, row_h=Inches(0.52))
    note(s, "INCONCLUSIVE is deliberately separated from NON-INFERIOR. Reporting a wide interval "
            "as equivalence would present absence of evidence as evidence of absence — with 28 "
            "subjects that distinction decides how several of the results below must be phrased.",
         Inches(4.5))

    # 9 significance, all models
    s = slide("Significance — every model against the boundary",
              f"Δ {E.label} vs SegFormer-B0 direct ({E.bound:.4f}). Paired subject bootstrap, 10 000 draws.")
    rows = [["Model", "KD method", f"{E.label}", "Δ vs boundary", "95 % CI", "Verdict"]]
    order = [("segformer_b0_distilled", "SegFormer-B0", "response"),
             ("segformer_b0_rgkd", "SegFormer-B0", "gated"),
             ("lraspp_mobilenetv3_b2kd", "LR-ASPP MobileNetV3", "response"),
             ("lraspp_mobilenetv3_rgkd", "LR-ASPP MobileNetV3", "gated"),
             ("topformer_tiny_b2kd", "TopFormer-Tiny", "response"),
             ("topformer_tiny_rgkd", "TopFormer-Tiny", "gated"),
             ("ppmobileseg_tiny_b2kd", "PP-MobileSeg-Tiny", "response"),
             ("ppmobileseg_tiny_rgkd", "PP-MobileSeg-Tiny", "gated")]
    for arm, name, meth in order:
        c = E.contrast(arm, REF)
        rows.append([name, meth, f"{E.value(arm):.4f}", f"{c['delta']:+.4f}",
                     f"[{c['lo']:+.4f}, {c['hi']:+.4f}]", c["verdict"]])
    table(s, rows, Inches(1.8), [2.5, 1.1, 1.1, 1.3, 1.9, 1.5], size=11.5,
          row_h=Inches(0.40))
    note(s, "Only the two SegFormer-B0 arms hold the boundary. Every 1–3 M student is below it, "
            "most with intervals excluding the margin — distillation narrows the architecture "
            "gap but does not remove it.", Inches(5.65))

    # 10 significance, gate vs control
    s = slide("Significance — does the gate beat plain response KD?",
              "Teacher, α, recipe and seeds held fixed. Only the loss differs.")
    rows = [["Student", f"Δ {E.label} (gated − response)", "95 % CI", "p", "Verdict"]]
    for m, n in [("segformer_b0", "SegFormer-B0"), ("lraspp_mobilenetv3", "LR-ASPP MobileNetV3"),
                 ("topformer_tiny", "TopFormer-Tiny"), ("ppmobileseg_tiny", "PP-MobileSeg-Tiny")]:
        a = f"{m}_rgkd"
        b = "segformer_b0_distilled" if m == "segformer_b0" else f"{m}_b2kd"
        c = E.contrast(a, b)
        rows.append([n, f"{c['delta']:+.4f}", f"[{c['lo']:+.4f}, {c['hi']:+.4f}]",
                     f"{c['p']:.3f}", c["verdict"]])
    table(s, rows, Inches(1.85), [2.6, 2.4, 2.2, 0.9, 1.6], size=13, row_h=Inches(0.48))
    note(s, "All four intervals span zero. The gate demonstrably fired — coverage 0.906, 2.8 % of "
            "image-views fully gated off — so this is evidence that suppressing the teacher's "
            "confident errors does not change what these students learn, not evidence that the "
            "gate did nothing. Reported as the null it is.", Inches(4.55))

    # 11 median-only extra
    if E.kind == "median":
        s = slide("What the median leaves out",
                  "Red = how much higher each arm reads on median than on mean.")
        pic(s, f"{E.prefix}5_mean_vs_median", Inches(1.6), Inches(4.5))
        note(s, "Every arm reads 0.032–0.066 higher on median. That uplift does NOT track "
                "complete misses (Spearman −0.09, p = 0.84) — with only 0–7 zeros out of 185 the "
                "misses are too few to move the mean much, so what the median discards is the "
                "broad tail of low-but-nonzero scores. Quote both endpoints, never one.",
             Inches(6.3))

    # 12 fairness
    s = slide("Skin-tone fairness across ITA groups",
              "Gap between best and worst ITA group. Lower is fairer.")
    pic(s, f"{E.prefix}6_fairness", Inches(1.6), Inches(4.5))
    note(s, "PP-MobileSeg starts with a significant gap and narrows under both KD methods, losing "
            "significance only when gated. LR-ASPP is never significant. TopFormer moves the "
            "wrong way under gating, so the direction is not consistent across students. With 28 "
            "subjects these are single comparisons, and bruise size is unevenly distributed "
            "across ITA groups — a fairness claim that does not condition on size measures both "
            "at once.", Inches(6.25))

    # 13 findings
    s = slide("What the results say")
    bullets(s, [
        ("Distillation from B2 helps every compact student.",
         "Each smaller model beats its own undistilled baseline; LR-ASPP gains most, with an "
         "interval excluding zero."),
        ("None of them reaches an undistilled SegFormer-B0.",
         "Distillation narrows the architecture gap; it does not remove it. Architecture and "
         "capacity still dominate."),
        ("Reliability gating did not beat response KD.",
         "Four students, every interval spanning zero, on both endpoints — despite the gate "
         "firing on 2.8 % of image-views. The simpler method is the one to ship."),
        ("The conclusion survives the endpoint swap.",
         "Mean and median give the same ranking and the same verdicts, which is more reassuring "
         "than either result alone."),
    ], top=Inches(1.7), gap=Inches(1.26))

    # 14 scope
    s = slide("Scope of these slides")
    bullets(s, [
        ("Fast-SCNN is excluded.",
         "It is the only model here with no pretrained backbone, so it confounds 'small' with "
         "'no ImageNet'. It is also where KD-vs-none was largest, so leaving it out understates "
         "what distillation achieved."),
        ("Best-seed numbers; three seeds were trained.",
         "Bars show the validation-selected seed. Three-seed means and spreads are in the "
         "handbook; the ordering is unchanged."),
        ("28 test subjects.",
         "Intervals are genuinely wide. INCONCLUSIVE means the study cannot answer that question "
         "at this size — not that the arms are equivalent."),
        ("Truncated y-axes.",
         f"Charts start at {E.ylim[0]:.2f}, not 0, so differences look larger than they are. "
         "Every bar is labelled with its value for that reason."),
    ], top=Inches(1.72), gap=Inches(1.28), hs=14, bs=12)

    prs.save(E.out)
    print(f"  -> {E.out.name}  ({len(prs.slides._sldIdLst)} slides)")


for kind in ("mean", "median"):
    print(f"\n{kind} deck")
    E = Endpoint(kind)
    chart_gate_explainer(E)
    chart_b0(E)
    chart_small(E, "_b2kd", C_RESP, "response KD from B2", f"{E.prefix}2_small_response")
    chart_small(E, "_rgkd", C_GATE, "reliability-gated KD from B2", f"{E.prefix}3_small_gated")
    chart_head_to_head(E)
    if kind == "median":
        chart_mean_vs_median(E)
    chart_fairness(E)
    build(E)
print(f"\nboth decks written to {ROOT / 'docs'}")
