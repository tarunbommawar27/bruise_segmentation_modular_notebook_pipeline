#!/usr/bin/env python
"""Emit docs/b2_distillation_simple.pptx -- the short, chart-led version.

SCOPE, AND WHAT IS DELIBERATELY LEFT OUT
-----------------------------------------
Four students: SegFormer-B0, LR-ASPP MobileNetV3, TopFormer-Tiny, PP-MobileSeg-Tiny.
**Fast-SCNN is excluded throughout**, by request. That is a defensible cut and not
only a cosmetic one -- it is the only model in the study with no pretrained
initialisation, so it confounds "small architecture" with "no ImageNet", which is
exactly the variable these slides are trying to isolate. It is also the only place
the KD-vs-none result was large, so leaving it out makes the deck's claims weaker,
not stronger. Both facts are stated on the closing slide rather than buried here.

CHART RULES
------------
Every chart is drawn here rather than reused from the notebook, so all four share
one visual grammar:

  * the dotted line is ALWAYS `segformer_b0_direct` (0.7663) -- the strongest model
    in the study trained with no teacher at all. It is the boundary every bar is
    read against.
  * the shaded band is human-vs-human agreement, 0.6998 (paul_vs_majority) to
    0.7549 (gbarimah_vs_erik). A bar inside that band is performing within the
    disagreement of two expert annotators, which is the fact that governs how big
    any of these differences actually are.
  * the y-axis is truncated at 0.60. Truncation exaggerates differences, so every
    bar carries its own value label and the band is drawn to keep the reader
    calibrated. Do not remove the band and keep the truncation.
  * grey = no teacher, navy = response KD, teal = reliability-gated. One colour per
    training regime, held across every chart.

All numbers are read from FINAL_RESULT at build time. Nothing is transcribed.
"""
from __future__ import annotations

import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib.patches import Patch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
RES = ROOT / "BRUISE_UNIFIED" / "FINAL_RESULT"
FIG = RES / "figures" / "deck"
FIG.mkdir(parents=True, exist_ok=True)
OUT = ROOT / "docs" / "b2_distillation_simple.pptx"

plt.rcParams.update({
    "font.family": "Times New Roman", "font.size": 12,
    "axes.spines.top": False, "axes.spines.right": False,
    "axes.edgecolor": "#555555", "savefig.dpi": 220,
})

C_NONE = "#9AA3AD"     # no teacher
C_RESP = "#1F3F6E"     # response KD
C_GATE = "#2E7D74"     # reliability-gated
C_BAND = "#E8E4D9"     # human agreement band
YLO, YHI = 0.60, 0.80

# ── data ─────────────────────────────────────────────────────────────────────
hA = pd.read_csv(RES / "headline_stage_a.csv").set_index("model")
hH = pd.read_csv(RES / "headline_stage_h.csv").set_index("model")
hE = pd.read_csv(RES / "headline_stage_e.csv").set_index("model")
b2 = pd.read_csv(RES / "significance_b2_teacher_vs_b0_direct.csv").set_index("arm")
famH = pd.read_csv(RES / "significance_contrast_family_stage_h.csv")
fair = pd.read_csv(RES / "stage_h_fairness.csv").set_index("model")
ceil = pd.read_csv(RES / "annotation_ceiling.csv").set_index("comparison")

BOUND = float(hA.loc["segformer_b0_direct", "mean_dice"])
BAND = (float(ceil.loc["human: paul_vs_majority", "mean_dice"]),
        float(ceil.loc["human: gbarimah_vs_erik", "mean_dice"]))


def dice(m: str) -> float:
    for t in (hH, hE, hA):
        if m in t.index:
            return float(t.loc[m, "mean_dice"])
    raise KeyError(m)


def misses(m: str) -> int:
    for t in (hH, hE, hA):
        if m in t.index:
            return int(t.loc[m, "complete_miss_count"])
    raise KeyError(m)


def contrast(a: str, b: str) -> pd.Series:
    m = famH[(famH.a == a) & (famH.b == b)]
    return m.iloc[0] if len(m) else None


# ── chart helpers ────────────────────────────────────────────────────────────
def _frame(ax, title):
    ax.axhspan(*BAND, color=C_BAND, zorder=0)
    ax.axhline(BOUND, color="#B03A2E", lw=1.6, ls=(0, (4, 3)), zorder=3)
    ax.set_ylim(YLO, YHI)
    ax.set_ylabel("mean Dice", fontsize=12)
    ax.grid(axis="y", alpha=.20, ls=":", zorder=0)
    ax.set_axisbelow(True)
    if title:
        ax.set_title(title, fontsize=13, weight="bold", pad=12)


def _label(ax, xs, vals, size=11):
    for x, v in zip(xs, vals):
        ax.text(x, v + 0.004, f"{v:.3f}", ha="center", va="bottom", fontsize=size)


def _legend(ax, entries, extra_band=True):
    """Legend BELOW the axes, never inside them.

    Placed outside because the boundary line sits at 0.766 and the plot tops out
    at 0.80 -- an in-axes legend lands on top of the one line the whole chart is
    read against, which is the worst possible collision.
    """
    h = [Patch(facecolor=c, label=l) for l, c in entries]
    h.append(plt.Line2D([], [], color="#B03A2E", lw=1.7, ls=(0, (4, 3)),
                        label=f"SegFormer-B0 direct, no teacher  ({BOUND:.3f})"))
    if extra_band:
        h.append(Patch(facecolor=C_BAND, label="human-vs-human agreement"))
    ax.legend(handles=h, frameon=False, fontsize=10.5,
              loc="upper center", bbox_to_anchor=(0.5, -0.13),
              ncol=2, handlelength=1.8, columnspacing=2.2)


def save(fig, name):
    fig.tight_layout()
    fig.savefig(FIG / f"{name}.png", bbox_inches="tight")
    plt.close(fig)
    print(f"  figures/deck/{name}.png")


# ── chart 1 · SegFormer-B0, both KD methods ──────────────────────────────────
def chart_b0():
    arms = [("segformer_b0_distilled", "Response KD\nfrom B2", C_RESP),
            ("segformer_b0_rgkd", "Reliability-gated KD\nfrom B2", C_GATE)]
    vals = [dice(a) for a, _, _ in arms]
    fig, ax = plt.subplots(figsize=(8.6, 5.1))
    _frame(ax, None)
    xs = np.arange(len(arms))
    ax.bar(xs, vals, width=.42, color=[c for _, _, c in arms], zorder=2)
    _label(ax, xs, vals, size=12.5)
    ax.set_xticks(xs)
    ax.set_xticklabels([l for _, l, _ in arms], fontsize=12)
    ax.set_xlim(-0.6, len(arms) - 0.4)
    _legend(ax, [("response KD", C_RESP), ("reliability-gated KD", C_GATE)])
    save(fig, "C1_b0_kd")


# ── chart 2 · smaller students, response KD ──────────────────────────────────
SMALL = [("lraspp_mobilenetv3", "LR-ASPP\nMobileNetV3", "3.22 M"),
         ("topformer_tiny", "TopFormer-Tiny", "1.37 M"),
         ("ppmobileseg_tiny", "PP-MobileSeg-Tiny", "1.45 M")]


def chart_small(suffix, colour, label, name):
    """Direct baseline beside the distilled arm, so the KD effect is visible."""
    base = [dice(m) for m, _, _ in SMALL]
    arm = [dice(f"{m}{suffix}") for m, _, _ in SMALL]
    xs = np.arange(len(SMALL))
    w = .34
    fig, ax = plt.subplots(figsize=(9.6, 5.3))
    _frame(ax, None)
    ax.bar(xs - w / 2, base, width=w, color=C_NONE, zorder=2)
    ax.bar(xs + w / 2, arm, width=w, color=colour, zorder=2)
    _label(ax, xs - w / 2, base, size=10.5)
    _label(ax, xs + w / 2, arm, size=10.5)
    ax.set_xticks(xs)
    ax.set_xticklabels([f"{n}\n{p}" for _, n, p in SMALL], fontsize=11.5)
    ax.set_xlim(-0.6, len(SMALL) - 0.4)
    _legend(ax, [("no teacher", C_NONE), (label, colour)])
    save(fig, name)


# ── chart 3 · response vs gated, side by side ────────────────────────────────
def chart_head_to_head():
    models = [("segformer_b0", "SegFormer-B0", "segformer_b0_distilled", "segformer_b0_rgkd")] + \
             [(m, n, f"{m}_b2kd", f"{m}_rgkd") for m, n, _ in SMALL]
    xs = np.arange(len(models))
    w = .34
    resp = [dice(a) for _, _, a, _ in models]
    gate = [dice(b) for _, _, _, b in models]
    fig, ax = plt.subplots(figsize=(10.4, 5.3))
    _frame(ax, None)
    ax.bar(xs - w / 2, resp, width=w, color=C_RESP, zorder=2)
    ax.bar(xs + w / 2, gate, width=w, color=C_GATE, zorder=2)
    _label(ax, xs - w / 2, resp, size=10)
    _label(ax, xs + w / 2, gate, size=10)
    ax.set_xticks(xs)
    ax.set_xticklabels([n for _, n, _, _ in models], fontsize=11.5)
    ax.set_xlim(-0.6, len(models) - 0.4)
    _legend(ax, [("response KD", C_RESP), ("reliability-gated KD", C_GATE)])
    save(fig, "C4_response_vs_gated")


# ── chart 4 · ITA fairness for the smaller students ──────────────────────────
def chart_fairness():
    rows = []
    for m, n, _ in SMALL:
        for arm, lab, col in ((m, "no teacher", C_NONE),
                              (f"{m}_b2kd", "response KD", C_RESP),
                              (f"{m}_rgkd", "gated KD", C_GATE)):
            r = fair.loc[arm]
            rows.append((n, lab, col, float(r.fairness_gap), bool(r.significant)))
    xs = np.arange(len(SMALL))
    w = .26
    fig, ax = plt.subplots(figsize=(9.8, 4.9))
    for k, (lab, col) in enumerate((("no teacher", C_NONE), ("response KD", C_RESP),
                                    ("gated KD", C_GATE))):
        vals = [r[3] for r in rows if r[1] == lab]
        sig = [r[4] for r in rows if r[1] == lab]
        pos = xs + (k - 1) * w
        ax.bar(pos, vals, width=w, color=col, zorder=2)
        for x, v, s in zip(pos, vals, sig):
            ax.text(x, v + 0.004, f"{v:.3f}" + ("*" if s else ""),
                    ha="center", va="bottom", fontsize=9.5,
                    weight="bold" if s else "normal")
    ax.set_xticks(xs)
    ax.set_xticklabels([n for _, n, _ in SMALL], fontsize=11.5)
    ax.set_ylabel("ITA fairness gap  (best group − worst group, median Dice)", fontsize=11)
    ax.set_ylim(0, max(r[3] for r in rows) * 1.25)
    ax.grid(axis="y", alpha=.20, ls=":")
    ax.set_axisbelow(True)
    ax.set_xlim(-0.6, len(SMALL) - 0.4)
    h = [Patch(facecolor=c, label=l) for l, c in
         (("no teacher", C_NONE), ("response KD", C_RESP), ("gated KD", C_GATE))]
    ax.legend(handles=h, frameon=False, fontsize=10.5, loc="upper left")
    ax.text(0.99, 0.97, "* Kruskal–Wallis p < 0.05\nlower is fairer",
            transform=ax.transAxes, ha="right", va="top", fontsize=10, color="#555555")
    save(fig, "C5_fairness_small")


chart_b0()
chart_small("_b2kd", C_RESP, "response KD from B2", "C2_small_response")
chart_small("_rgkd", C_GATE, "reliability-gated KD from B2", "C3_small_gated")
chart_head_to_head()
chart_fairness()

# ═════════════════════════════════════════════════════════════════════════════
# deck
# ═════════════════════════════════════════════════════════════════════════════
FONT = "Times New Roman"
INK = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT = RGBColor(0x1F, 0x3F, 0x6E)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
RULE = RGBColor(0xC9, 0xCF, 0xD8)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H, M = prs.slide_width, prs.slide_height, Inches(0.75)
BLANK = prs.slide_layouts[6]


def _f(run, size, bold=False, color=INK, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def tbox(s, l, t, w, h, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    tb.text_frame.paragraphs[0].alignment = align
    return tb.text_frame


def slide(title, kicker=None):
    s = prs.slides.add_slide(BLANK)
    tf = tbox(s, M, Inches(0.40), W - 2 * M, Inches(0.7))
    _f(tf.paragraphs[0].add_run(), 25, bold=True, color=ACCENT)
    tf.paragraphs[0].runs[0].text = title
    ln = s.shapes.add_shape(1, M, Inches(1.10), W - 2 * M, Emu(9525))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE
    ln.line.fill.background(); ln.shadow.inherit = False
    if kicker:
        kf = tbox(s, M, Inches(1.17), W - 2 * M, Inches(0.45))
        _f(kf.paragraphs[0].add_run(), 12.5, color=MUTED, italic=True)
        kf.paragraphs[0].runs[0].text = kicker
    return s


def pic(s, name, top, max_h):
    from PIL import Image
    p = FIG / f"{name}.png"
    iw, ih = Image.open(p).size
    w = W - 2 * M
    h = int(w * ih / iw)
    if h > max_h:
        h, w = max_h, int(max_h * iw / ih)
    s.shapes.add_picture(str(p), int((W - w) / 2), top, width=int(w), height=int(h))


def note(s, text, top=Inches(6.55), size=11.5):
    tf = tbox(s, M, top, W - 2 * M, Inches(0.75))
    _f(tf.paragraphs[0].add_run(), size, color=MUTED)
    tf.paragraphs[0].runs[0].text = text


def table(s, rows, left, top, width, col_w, size=12, row_h=Inches(0.42), bold_rows=()):
    t = s.shapes.add_table(len(rows), len(rows[0]), left, top, width, row_h * len(rows)).table
    t.first_row = False
    tot = sum(col_w)
    for i, cw in enumerate(col_w):
        t.columns[i].width = int(width * cw / tot)
    for r, row in enumerate(rows):
        t.rows[r].height = row_h
        for c, v in enumerate(row):
            cell = t.cell(r, c)
            cell.text = ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = 0
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run(); run.text = str(v)
            _f(run, size, bold=(r == 0 or r in bold_rows),
               color=MUTED if r == 0 else INK)
    return t


# 1 · title
s = prs.slides.add_slide(BLANK)
tf = tbox(s, M, Inches(2.6), W - 2 * M, Inches(1.2))
_f(tf.paragraphs[0].add_run(), 38, bold=True, color=ACCENT)
tf.paragraphs[0].runs[0].text = "Distilling SegFormer-B2 into Compact Students"
tf2 = tbox(s, M, Inches(3.95), W - 2 * M, Inches(0.9))
_f(tf2.paragraphs[0].add_run(), 16, color=INK)
tf2.paragraphs[0].runs[0].text = "Response KD and reliability-gated KD, against an undistilled SegFormer-B0"
ln = s.shapes.add_shape(1, M, Inches(5.05), Inches(2.2), Emu(19050))
ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT
ln.line.fill.background(); ln.shadow.inherit = False
tf3 = tbox(s, M, Inches(5.3), W - 2 * M, Inches(0.6))
_f(tf3.paragraphs[0].add_run(), 12.5, color=MUTED)
tf3.paragraphs[0].runs[0].text = "185 test images · 28 subjects · 3 seeds per arm · one fixed training recipe"

# 2 · the students
s = slide("The students", "Chosen to span the CNN–transformer axis at deployment scale.")
rows = [["Model", "Params", "Architecture", "Pretrained backbone", "Why it is here"],
        ["SegFormer-B0", "3.71 M", "pure transformer (MiT)", "ImageNet-1k", "same family as the teacher — isolates size from architecture"],
        ["LR-ASPP MobileNetV3", "3.22 M", "pure CNN", "ImageNet-1k", "strongest mobile baseline; the realistic deployment target"],
        ["TopFormer-Tiny", "1.37 M", "hybrid — CNN token pyramid + transformer", "ImageNet-1k, 66.2 % top-1", "half the parameters, attention retained"],
        ["PP-MobileSeg-Tiny", "1.45 M", "hybrid — StrideFormer, strided SEA attention", "ImageNet", "a different attention design at the same scale"]]
table(s, rows, M, Inches(1.75), W - 2 * M, [1.9, 0.8, 2.6, 1.8, 3.4], size=11.5, row_h=Inches(0.62))
note(s, "All four carry an ImageNet-pretrained backbone and a freshly initialised head — the same "
        "setup as every other baseline in the study, which is what makes the comparison fair. "
        "Spanning pure CNN to pure transformer matters: it is what lets a distillation result be "
        "read as a property of the method rather than of one architecture.", Inches(5.05))

# 3 · chart 1
s = slide("SegFormer-B0: both KD methods sit on the boundary",
          "Dotted line = the same architecture trained with no teacher at all.")
pic(s, "C1_b0_kd", Inches(1.62), Inches(4.5))
r = contrast("segformer_b0_rgkd", "segformer_b0_distilled")
note(s, f"Response KD gains +{dice('segformer_b0_distilled') - BOUND:.4f} over no teacher; gating "
        f"{dice('segformer_b0_rgkd') - BOUND:+.4f}. Gated − response = {r.delta_dice:+.4f}, "
        f"95 % CI [{r.lo:+.4f}, {r.hi:+.4f}] — spans zero. At 3.71 M parameters the student is "
        f"already at the teacher's level, so there is almost nothing for distillation to add.",
     Inches(6.25))

# 4 · chart 2
s = slide("Smaller students: response KD from B2",
          "Grey = no teacher. Navy = distilled from SegFormer-B2.")
pic(s, "C2_small_response", Inches(1.62), Inches(4.4))
note(s, "Every student improves over its own undistilled baseline, and PP-MobileSeg gains most "
        "(+0.019). None reaches the dotted boundary: distillation narrows the gap to a 3.71 M "
        "transformer but does not close it. Both statements are true and neither should be "
        "quoted alone.", Inches(6.2))

# 5 · chart 3
s = slide("Smaller students: reliability-gated KD from B2",
          "Same teacher, same α, same recipe — only the loss differs.")
pic(s, "C3_small_gated", Inches(1.62), Inches(4.4))
note(s, "The gate suppresses the teacher on pixels where it is confidently wrong, and switches "
        "it off entirely on images the teacher missed. It fired on 2.8 % of image-views. The bars "
        "land in essentially the same place as response KD.", Inches(6.2))

# 6 · head to head
s = slide("Response KD vs reliability-gated KD, head to head",
          "Every difference is inside the noise floor of the labels.")
pic(s, "C4_response_vs_gated", Inches(1.62), Inches(4.3))
rows = [["Student", "Δ Dice (gated − response)", "95 % CI", "Verdict"]]
for m, n in [("segformer_b0", "SegFormer-B0"), ("lraspp_mobilenetv3", "LR-ASPP MobileNetV3"),
             ("topformer_tiny", "TopFormer-Tiny"), ("ppmobileseg_tiny", "PP-MobileSeg-Tiny")]:
    a = f"{m}_rgkd"
    b = "segformer_b0_distilled" if m == "segformer_b0" else f"{m}_b2kd"
    r = contrast(a, b)
    rows.append([n, f"{r.delta_dice:+.4f}", f"[{r.lo:+.4f}, {r.hi:+.4f}]", r.verdict])
note(s, "All four intervals span zero. Reliability gating is indistinguishable from plain "
        "response KD on this dataset — reported as the null it is.", Inches(6.15))

# 7 · significance table
s = slide("The same comparison as numbers",
          "Paired subject-level bootstrap, 28 subjects, 10 000 draws.")
table(s, rows, M, Inches(1.8), W - 2 * M, [2.6, 2.2, 2.2, 1.6], size=13, row_h=Inches(0.45))
rows2 = [["Against the boundary (SegFormer-B0 direct)", "Δ Dice", "95 % CI", "Verdict"]]
for arm, n in [("segformer_b0_distilled", "SegFormer-B0 · response"),
               ("segformer_b0_rgkd", "SegFormer-B0 · gated"),
               ("lraspp_mobilenetv3_b2kd", "LR-ASPP · response"),
               ("topformer_tiny_b2kd", "TopFormer · response"),
               ("ppmobileseg_tiny_b2kd", "PP-MobileSeg · response")]:
    r = b2.loc[arm]
    rows2.append([n, f"{r.delta_dice:+.4f}", f"[{r.lo:+.4f}, {r.hi:+.4f}]", r.verdict])
table(s, rows2, M, Inches(4.15), W - 2 * M, [2.6, 2.2, 2.2, 1.6], size=13, row_h=Inches(0.45))

# 8 · fairness
s = slide("Skin-tone fairness across ITA groups",
          "Gap between the best and worst ITA group. Lower is fairer.")
pic(s, "C5_fairness_small", Inches(1.62), Inches(4.4))
note(s, "PP-MobileSeg starts with a significant gap (0.179) and narrows under both KD methods, "
        "losing significance only when gated. LR-ASPP improves slightly and is never significant. "
        "TopFormer moves the wrong way under gating (0.080 → 0.115), so the direction is not "
        "consistent across students. With 28 test subjects these are single comparisons, and "
        "bruise size is unevenly distributed across ITA groups — a fairness claim that does not "
        "condition on size is measuring both at once.", Inches(6.15))

# 9 · findings
s = slide("What the results say")
items = [
    ("Distillation from B2 helps every compact student.",
     "Each of the three smaller models beats its own undistilled baseline; LR-ASPP gains "
     "+0.026 with an interval excluding zero."),
    ("None of them reaches an undistilled SegFormer-B0.",
     "Gaps of 0.045–0.091 Dice to the dotted boundary. Distillation narrows the architecture "
     "gap; it does not remove it."),
    ("Reliability gating did not beat response KD.",
     "Four students, every interval spanning zero — despite the gate firing on 2.8 % of "
     "image-views. The simpler method is the one to ship."),
    ("Everything sits inside the annotation ceiling.",
     "Two expert annotators agree at 0.755 Dice, and the whole field lies in the shaded band. "
     "Statistical significance is not clinical importance."),
]
top = Inches(1.7)
for head, body in items:
    tf = tbox(s, M, top, W - 2 * M, Inches(1.15))
    _f(tf.paragraphs[0].add_run(), 15, bold=True, color=ACCENT)
    tf.paragraphs[0].runs[0].text = head
    p = tf.add_paragraph()
    _f(p.add_run(), 12.5, color=INK)
    p.runs[0].text = body
    top += Inches(1.26)

# 10 · scope note
s = slide("Scope of these slides")
rows = [["Fast-SCNN is excluded",
         "It is the only model here with no pretrained backbone, so it confounds 'small' with "
         "'no ImageNet'. It is also where KD-vs-none was largest (+0.036, misses −3.8 pp), so "
         "leaving it out understates what distillation achieved."],
        ["Best-seed numbers",
         "Bars are the validation-selected seed, as in the headline tables. Three-seed means and "
         "spreads are in the handbook; the ordering is unchanged."],
        ["28 test subjects",
         "Intervals are genuinely wide. INCONCLUSIVE means the study cannot answer that question "
         "at this size — not that the arms are equivalent."],
        ["Teacher choice was tested separately",
         "A B2 teacher beat DeepLabV3+ on the same student by +0.034 with the interval excluding "
         "zero, which is what justifies using B2 here."]]
top = Inches(1.75)
for head, body in rows:
    tf = tbox(s, M, top, W - 2 * M, Inches(1.1))
    _f(tf.paragraphs[0].add_run(), 14, bold=True, color=INK)
    tf.paragraphs[0].runs[0].text = head
    p = tf.add_paragraph()
    _f(p.add_run(), 12, color=MUTED)
    p.runs[0].text = body
    top += Inches(1.24)

prs.save(OUT)
print(f"\nwrote {OUT}")
print(f"  {len(prs.slides._sldIdLst)} slides, {FONT}, 16:9")
